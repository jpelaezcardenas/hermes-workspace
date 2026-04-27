#!/usr/bin/env python3
"""Run ppt-maker active runtime golden regressions."""
from __future__ import annotations

import argparse
import json
import shutil
import subprocess
import sys
import time
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[2]
WRAPPER = ROOT / "bin" / "ppt-maker"
PPT = ROOT / "legacy" / "ppt-maker"
BASE_CONFIG = PPT / "configs" / "dalp-overview-10-slides.json"


def run_case(name: str, cmd: list[str], cwd: Path | None = None) -> dict[str, Any]:
    start = time.time()
    proc = subprocess.run(cmd, cwd=str(cwd or ROOT), text=True, capture_output=True)
    return {
        "name": name,
        "cmd": cmd,
        "exit": proc.returncode,
        "ok": proc.returncode == 0,
        "duration_seconds": round(time.time() - start, 3),
        "stdout_tail": proc.stdout[-4000:],
        "stderr_tail": proc.stderr[-4000:],
    }


def pptx_gate(path: Path, expected_slides: int) -> dict[str, Any]:
    result: dict[str, Any] = {"path": str(path), "exists": path.exists(), "ok": False}
    if not path.exists():
        return result
    try:
        with zipfile.ZipFile(path) as zf:
            bad = zf.testzip()
            names = zf.namelist()
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        full_text = "\n".join(texts)
        result.update({
            "zip_ok": bad is None,
            "has_presentation_xml": "ppt/presentation.xml" in names,
            "slide_count": len(prs.slides),
            "text_chars": len(full_text),
            "has_dalp_text": "DALP" in full_text or "Digital Asset Lifecycle Platform" in full_text,
            "placeholder_hits": [s for s in ["Lorem", "ipsum", "[VARIABLE", "TODO"] if s.lower() in full_text.lower()],
        })
        result["ok"] = bool(result["zip_ok"] and result["has_presentation_xml"] and result["slide_count"] == expected_slides and result["text_chars"] > 300 and result["has_dalp_text"] and not result["placeholder_hits"])
    except Exception as exc:  # noqa: BLE001
        result["error"] = repr(exc)
    return result


def pdf_gate(path: Path, expected_pages: int) -> dict[str, Any]:
    result: dict[str, Any] = {"path": str(path), "exists": path.exists(), "ok": False}
    if not path.exists():
        return result
    info = run_case("pdfinfo", ["pdfinfo", str(path)])
    pages = None
    for line in info["stdout_tail"].splitlines():
        if line.startswith("Pages:"):
            pages = int(line.split(":", 1)[1].strip())
    result.update({"pdfinfo_ok": info["ok"], "pages": pages})
    result["ok"] = bool(info["ok"] and pages == expected_pages)
    return result


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Run ppt-maker active runtime golden regressions")
    parser.add_argument("--json", action="store_true")
    parser.add_argument("--out-dir", default=None)
    ns = parser.parse_args(argv)

    stamp = time.strftime("%Y%m%dT%H%M%SZ", time.gmtime())
    out_dir = Path(ns.out_dir).resolve() if ns.out_dir else PPT / "output" / f"golden-{stamp}"
    out_dir.mkdir(parents=True, exist_ok=True)

    config = json.loads(BASE_CONFIG.read_text(encoding="utf-8"))
    config["project"] = "ppt-maker-golden-dalp-overview"
    config["output"] = str((out_dir / "ppt-maker-golden-dalp-overview.pptx").relative_to(PPT))
    config_path = out_dir / "ppt-maker-golden-dalp-overview.config.json"
    config_path.write_text(json.dumps(config, indent=2), encoding="utf-8")
    pptx = out_dir / "ppt-maker-golden-dalp-overview.pptx"
    pdf = out_dir / "ppt-maker-golden-dalp-overview.pdf"

    cases: list[dict[str, Any]] = []
    cases.append(run_case("smoke", [str(WRAPPER), "smoke"]))
    cases.append(run_case("build_config", [str(WRAPPER), "build-config", str(config_path)]))
    cases.append(run_case("inspect", [str(WRAPPER), "inspect", str(pptx), "--json"]))
    cases.append(run_case("validate", [str(WRAPPER), "validate", str(pptx)]))
    cases.append(run_case("render_pdf", [str(WRAPPER), "render-pdf", str(pptx), str(out_dir)]))
    if pdf.exists():
        cases.append(run_case("render_first_slide_image", ["pdftoppm", "-png", "-r", "100", "-f", "1", "-l", "1", str(pdf), str(out_dir / "slide")]))
    cases.append(run_case("visual_regression", [sys.executable, str(ROOT / "active" / "runtime-scripts" / "run_visual_regression.py"), "--json", "--out-dir", str(out_dir / "visual-regression")]))

    gates = [pptx_gate(pptx, 10), pdf_gate(pdf, 10)]
    unexpected = [c for c in cases if not c["ok"]]
    failed_gates = [g for g in gates if not g["ok"]]
    payload = {
        "agent": "ppt-maker",
        "root": str(ROOT),
        "artifact_dir": str(out_dir),
        "pptx": str(pptx),
        "pdf": str(pdf),
        "case_count": len(cases),
        "passed_count": sum(1 for c in cases if c["ok"]),
        "gate_count": len(gates),
        "gate_passed_count": sum(1 for g in gates if g["ok"]),
        "unexpected": unexpected,
        "failed_gates": failed_gates,
        "cases": cases,
        "gates": gates,
        "ok": not unexpected and not failed_gates,
    }
    (out_dir / "golden-regressions.json").write_text(json.dumps(payload, indent=2), encoding="utf-8")
    if ns.json:
        print(json.dumps(payload, indent=2))
    else:
        print(f"ppt-maker golden regressions: {payload['passed_count']}/{payload['case_count']} cases and {payload['gate_passed_count']}/{payload['gate_count']} gates passed")
        print(f"artifacts: {out_dir}")
    return 0 if payload["ok"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
