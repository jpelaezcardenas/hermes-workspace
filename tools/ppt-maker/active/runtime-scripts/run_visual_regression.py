#!/usr/bin/env python3
"""Run ppt-maker visual/client-artifact regression with screenshot-heavy deck and checker review."""
from __future__ import annotations

import argparse
import json
import shutil
import subprocess
import time
import zipfile
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[2]
WRAPPER = ROOT / "bin" / "ppt-maker"
PPT = ROOT / "legacy" / "ppt-maker"
BASE_CONFIG = PPT / "configs" / "2026-03-24-dalp-platform-walkthrough.json"


def run_case(name: str, cmd: list[str], cwd: Path | None = None) -> dict[str, Any]:
    start = time.time()
    proc = subprocess.run(cmd, cwd=str(cwd or ROOT), text=True, capture_output=True)
    return {
        "name": name,
        "cmd": cmd,
        "exit": proc.returncode,
        "ok": proc.returncode == 0,
        "duration_seconds": round(time.time() - start, 3),
        "stdout_tail": proc.stdout[-4000:],
        "stderr_tail": proc.stderr[-4000:],
    }


def pptx_gate(path: Path, min_slides: int, min_images: int) -> dict[str, Any]:
    result: dict[str, Any] = {"path": str(path), "exists": path.exists(), "ok": False}
    if not path.exists():
        return result
    try:
        with zipfile.ZipFile(path) as zf:
            bad = zf.testzip()
            names = zf.namelist()
            media = [n for n in names if n.startswith("ppt/media/")]
        prs = Presentation(str(path))
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        full_text = "\n".join(texts)
        result.update({
            "zip_ok": bad is None,
            "has_presentation_xml": "ppt/presentation.xml" in names,
            "slide_count": len(prs.slides),
            "media_count": len(media),
            "text_chars": len(full_text),
            "has_dalp_text": "DALP" in full_text or "Digital Asset Lifecycle Platform" in full_text,
            "placeholder_hits": [s for s in ["Lorem", "ipsum", "[VARIABLE", "TODO"] if s.lower() in full_text.lower()],
        })
        result["ok"] = bool(
            result["zip_ok"]
            and result["has_presentation_xml"]
            and result["slide_count"] >= min_slides
            and result["media_count"] >= min_images
            and result["text_chars"] > 1200
            and result["has_dalp_text"]
            and not result["placeholder_hits"]
        )
    except Exception as exc:  # noqa: BLE001
        result["error"] = repr(exc)
    return result


def pdf_gate(path: Path, min_pages: int) -> dict[str, Any]:
    result: dict[str, Any] = {"path": str(path), "exists": path.exists(), "ok": False}
    if not path.exists():
        return result
    info = run_case("pdfinfo", ["pdfinfo", str(path)])
    pages = None
    for line in info["stdout_tail"].splitlines():
        if line.startswith("Pages:"):
            pages = int(line.split(":", 1)[1].strip())
    result.update({"pdfinfo_ok": info["ok"], "pages": pages})
    result["ok"] = bool(info["ok"] and pages and pages >= min_pages)
    return result



def make_fake_screenshot(path: Path, title: str, idx: int) -> None:
    """Create a deterministic dashboard-like image for screenshot placeholder regression."""
    path.parent.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGB", (1600, 900), "#f8fafc")
    draw = ImageDraw.Draw(img)
    # Browser chrome
    draw.rounded_rectangle((20, 20, 1580, 880), radius=24, fill="white", outline="#cbd5e1", width=3)
    draw.rectangle((20, 20, 1580, 88), fill="#0f172a")
    for x, c in [(52, "#ef4444"), (84, "#f59e0b"), (116, "#22c55e")]:
        draw.ellipse((x, 45, x + 16, 61), fill=c)
    draw.rounded_rectangle((170, 38, 760, 70), radius=10, fill="#1e293b")
    draw.text((190, 45), "dalp.settlemint.com / workspace", fill="#e2e8f0")
    # Sidebar
    draw.rectangle((40, 110, 300, 850), fill="#111827")
    draw.text((70, 140), "SettleMint DALP", fill="#ffffff")
    nav = ["Dashboard", "Assets", "Investors", "Compliance", "Operations", "Audit"]
    for i, item in enumerate(nav):
        y = 205 + i * 70
        fill = "#2563eb" if i == idx % len(nav) else "#1f2937"
        draw.rounded_rectangle((65, y, 270, y + 42), radius=12, fill=fill)
        draw.text((88, y + 12), item, fill="#ffffff")
    # Main content
    draw.text((340, 130), title[:54], fill="#0f172a")
    colors = ["#dbeafe", "#dcfce7", "#fef3c7", "#fee2e2"]
    labels = ["AUM", "Pending approvals", "Assets live", "Policy checks"]
    for i, label in enumerate(labels):
        x = 340 + (i % 2) * 370
        y = 200 + (i // 2) * 150
        draw.rounded_rectangle((x, y, x + 330, y + 110), radius=18, fill=colors[i], outline="#cbd5e1")
        draw.text((x + 25, y + 25), label, fill="#334155")
        draw.text((x + 25, y + 60), str((idx + 2) * (i + 3) * 17), fill="#0f172a")
    # Chart/table
    draw.rounded_rectangle((340, 540, 1500, 820), radius=18, fill="#ffffff", outline="#cbd5e1")
    for i in range(8):
        x = 380 + i * 130
        h = 40 + ((i + idx) % 6) * 28
        draw.rectangle((x, 790 - h, x + 70, 790), fill="#2563eb")
    for i in range(5):
        y = 570 + i * 42
        draw.line((930, y, 1460, y), fill="#e2e8f0", width=2)
        draw.text((950, y + 10), f"Transaction event {i+1}", fill="#475569")
    img.save(path)

def make_montage(images: list[Path], out_path: Path, cols: int = 4) -> dict[str, Any]:
    thumbs: list[Image.Image] = []
    for img_path in images:
        img = Image.open(img_path).convert("RGB")
        img.thumbnail((360, 205))
        canvas = Image.new("RGB", (380, 235), "white")
        canvas.paste(img, ((380 - img.width) // 2, 20))
        draw = ImageDraw.Draw(canvas)
        draw.text((8, 5), img_path.stem, fill=(0, 0, 0))
        thumbs.append(canvas)
    if not thumbs:
        return {"ok": False, "image_count": 0}
    rows = (len(thumbs) + cols - 1) // cols
    montage = Image.new("RGB", (cols * 380, rows * 235), "#f3f4f6")
    for i, thumb in enumerate(thumbs):
        montage.paste(thumb, ((i % cols) * 380, (i // cols) * 235))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    montage.save(out_path)
    return {"ok": out_path.exists(), "image_count": len(thumbs), "path": str(out_path), "size": montage.size}


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Run screenshot-heavy ppt-maker visual regression")
    parser.add_argument("--json", action="store_true")
    parser.add_argument("--out-dir", default=None)
    ns = parser.parse_args(argv)

    stamp = time.strftime("%Y%m%dT%H%M%SZ", time.gmtime())
    out_dir = Path(ns.out_dir).resolve() if ns.out_dir else PPT / "output" / f"visual-{stamp}"
    out_dir.mkdir(parents=True, exist_ok=True)

    config = json.loads(BASE_CONFIG.read_text(encoding="utf-8"))
    # Keep the regression substantial but fast enough: cover + agenda + 8 screenshot slides + close.
    selected = config["slides"][:18] + [{"bank_file": "slide-024-thank-you.pptx"}]
    fake_dir = out_dir / "synthetic-screenshots"
    shot_idx = 0
    for slide in selected:
        if slide.get("image_path"):
            shot_idx += 1
            shot = fake_dir / f"dalp-screenshot-{shot_idx:02d}.png"
            make_fake_screenshot(shot, slide.get("title") or "DALP workspace", shot_idx)
            slide["image_path"] = str(shot)
    config["project"] = "ppt-maker-visual-dalp-walkthrough"
    config["description"] = "Visual regression fixture for ppt-maker screenshot and review lanes."
    config["slides"] = selected
    config["output"] = str((out_dir / "ppt-maker-visual-dalp-walkthrough.pptx").relative_to(PPT))
    config_path = out_dir / "ppt-maker-visual-dalp-walkthrough.config.json"
    config_path.write_text(json.dumps(config, indent=2), encoding="utf-8")

    pptx = out_dir / "ppt-maker-visual-dalp-walkthrough.pptx"
    pdf = out_dir / "ppt-maker-visual-dalp-walkthrough.pdf"
    review_md = out_dir / "ppt-maker-visual-review.md"
    review_json = out_dir / "ppt-maker-visual-review.json"
    montage = out_dir / "ppt-maker-visual-montage.png"

    cases: list[dict[str, Any]] = []
    cases.append(run_case("build_visual_config", [str(WRAPPER), "build-config", str(config_path)]))
    cases.append(run_case("inspect_visual", [str(WRAPPER), "inspect", str(pptx), "--json"]))
    cases.append(run_case("validate_visual", [str(WRAPPER), "validate", str(pptx)]))
    cases.append(run_case("generate_review", [str(WRAPPER), "generate-review", str(pptx), "--output-md", str(review_md), "--output-json", str(review_json)]))
    cases.append(run_case("render_pdf", [str(WRAPPER), "render-pdf", str(pptx), str(out_dir)]))

    rendered_images: list[Path] = []
    if pdf.exists():
        prefix = out_dir / "slide"
        cases.append(run_case("render_slide_images", ["pdftoppm", "-png", "-r", "90", str(pdf), str(prefix)]))
        rendered_images = sorted(out_dir.glob("slide-*.png"))
    montage_gate = make_montage(rendered_images, montage)

    gates = [
        pptx_gate(pptx, min_slides=12, min_images=8),
        pdf_gate(pdf, min_pages=12),
        {"name": "review_md", "path": str(review_md), "exists": review_md.exists(), "ok": review_md.exists() and review_md.stat().st_size > 1500},
        {"name": "review_json", "path": str(review_json), "exists": review_json.exists(), "ok": review_json.exists() and review_json.stat().st_size > 3000},
        {"name": "montage", **montage_gate},
    ]
    unexpected = [c for c in cases if not c["ok"]]
    failed_gates = [g for g in gates if not g.get("ok")]
    payload = {
        "agent": "ppt-maker",
        "root": str(ROOT),
        "artifact_dir": str(out_dir),
        "pptx": str(pptx),
        "pdf": str(pdf),
        "review_md": str(review_md),
        "review_json": str(review_json),
        "montage": str(montage),
        "case_count": len(cases),
        "passed_count": sum(1 for c in cases if c["ok"]),
        "gate_count": len(gates),
        "gate_passed_count": sum(1 for g in gates if g.get("ok")),
        "unexpected": unexpected,
        "failed_gates": failed_gates,
        "cases": cases,
        "gates": gates,
        "ok": not unexpected and not failed_gates,
    }
    (out_dir / "visual-regression.json").write_text(json.dumps(payload, indent=2), encoding="utf-8")
    latest = ROOT / "verification" / "visual-regression-latest.json"
    latest.parent.mkdir(parents=True, exist_ok=True)
    latest.write_text(json.dumps(payload, indent=2), encoding="utf-8")
    if ns.json:
        print(json.dumps(payload, indent=2))
    else:
        print(f"ppt-maker visual regression: {payload['passed_count']}/{payload['case_count']} cases and {payload['gate_passed_count']}/{payload['gate_count']} gates passed")
        print(f"artifacts: {out_dir}")
    return 0 if payload["ok"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
