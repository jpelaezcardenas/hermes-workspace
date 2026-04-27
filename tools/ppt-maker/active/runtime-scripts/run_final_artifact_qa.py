#!/usr/bin/env python3
"""Run ppt-maker client-ready final artifact QA for a deck and review bundle."""
from __future__ import annotations

import argparse
import json
import subprocess
import tempfile
import time
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[2]
WRAPPER = ROOT / "bin" / "ppt-maker"


def run(cmd: list[str], cwd: Path | None = None) -> dict[str, Any]:
    start = time.time()
    proc = subprocess.run(cmd, cwd=str(cwd or ROOT), text=True, capture_output=True)
    return {
        "cmd": cmd,
        "exit": proc.returncode,
        "ok": proc.returncode == 0,
        "duration_seconds": round(time.time() - start, 3),
        "stdout_tail": proc.stdout[-5000:],
        "stderr_tail": proc.stderr[-5000:],
    }


def pptx_gate(path: Path) -> dict[str, Any]:
    result: dict[str, Any] = {"path": str(path), "exists": path.exists(), "ok": False}
    if not path.exists():
        return result
    try:
        with zipfile.ZipFile(path) as zf:
            bad = zf.testzip()
            names = zf.namelist()
        prs = Presentation(str(path))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
        joined = "\n".join(text)
        result.update({
            "zip_ok": bad is None,
            "slide_count": len(prs.slides),
            "has_presentation_xml": "ppt/presentation.xml" in names,
            "text_chars": len(joined),
        })
        result["ok"] = bool(result["zip_ok"] and result["has_presentation_xml"] and result["slide_count"] >= 1 and result["text_chars"] >= 200)
    except Exception as exc:  # noqa: BLE001
        result["error"] = repr(exc)
    return result


def docx_gate(path: Path) -> dict[str, Any]:
    result: dict[str, Any] = {"path": str(path), "exists": path.exists(), "ok": False}
    if not path.exists():
        return result
    try:
        with zipfile.ZipFile(path) as zf:
            bad = zf.testzip()
            names = zf.namelist()
        result.update({
            "zip_ok": bad is None,
            "has_document_xml": "word/document.xml" in names,
            "has_styles": "word/styles.xml" in names,
        })
        result["ok"] = bool(result["zip_ok"] and result["has_document_xml"] and result["has_styles"])
    except Exception as exc:  # noqa: BLE001
        result["error"] = repr(exc)
    return result


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Run final artifact QA for ppt-maker outputs")
    parser.add_argument("--json", action="store_true")
    parser.add_argument("deck")
    parser.add_argument("--out-dir", default=None)
    ns = parser.parse_args(argv)

    deck = Path(ns.deck).resolve()
    out_dir = Path(ns.out_dir).resolve() if ns.out_dir else deck.parent / "final-qa"
    out_dir.mkdir(parents=True, exist_ok=True)

    review_md = out_dir / f"{deck.stem}-review.md"
    review_json = out_dir / f"{deck.stem}-review.json"
    review_docx = out_dir / f"{deck.stem}-review.docx"

    cases = [
        run([str(WRAPPER), "inspect", str(deck), "--json"]),
        run([str(WRAPPER), "validate", str(deck)]),
        run([str(WRAPPER), "generate-review", str(deck), "--output-md", str(review_md), "--output-json", str(review_json)]),
        run([str(WRAPPER), "export-review-docx", str(review_md), str(review_docx)]),
    ]

    gates = {
        "deck": pptx_gate(deck),
        "review_md": {"path": str(review_md), "exists": review_md.exists(), "ok": review_md.exists() and review_md.stat().st_size > 1000},
        "review_json": {"path": str(review_json), "exists": review_json.exists(), "ok": review_json.exists() and review_json.stat().st_size > 1000},
        "review_docx": docx_gate(review_docx),
    }

    payload = {
        "agent": "ppt-maker",
        "deck": str(deck),
        "artifact_dir": str(out_dir),
        "review_md": str(review_md),
        "review_json": str(review_json),
        "review_docx": str(review_docx),
        "case_count": len(cases),
        "passed_count": sum(1 for c in cases if c["ok"]),
        "gate_count": len(gates),
        "gate_passed_count": sum(1 for g in gates.values() if g.get("ok")),
        "cases": cases,
        "gates": gates,
        "ok": all(c["ok"] for c in cases) and all(g.get("ok") for g in gates.values()),
    }

    (out_dir / "final-artifact-qa.json").write_text(json.dumps(payload, indent=2), encoding="utf-8")
    if ns.json:
        print(json.dumps(payload, indent=2))
    else:
        print(f"ppt-maker final artifact QA: {payload['passed_count']}/{payload['case_count']} cases and {payload['gate_passed_count']}/{payload['gate_count']} gates passed")
        print(f"artifacts: {out_dir}")
    return 0 if payload["ok"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
