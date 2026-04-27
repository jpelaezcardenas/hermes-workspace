#!/usr/bin/env python3
"""Reusable image and diagram helpers for rich, editable PPTX slides."""
from __future__ import annotations

import hashlib
import json
import math
from pathlib import Path
from typing import Any

from PIL import Image, ImageColor, ImageDraw
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
LOGO_DIR = ROOT / "ppt-maker" / "assets" / "logos"
LOGO_REGISTRY_PATH = LOGO_DIR / "logo_registry.json"
ICON_REGISTRY_PATH = ROOT / "ppt-maker" / "assets" / "icons" / "icon_registry.json"

BRAND_COLORS: dict[str, str] = {
    "navy": "1B2B65",
    "blue": "3B82F6",
    "green": "10B981",
    "amber": "F59E0B",
    "red": "EF4444",
    "purple": "8B5CF6",
    "light_bg": "F8FAFC",
    "gray": "64748B",
    "mid_gray": "CBD5E1",
    "dark": "1E293B",
    "white": "FFFFFF",
    "black": "0F172A",
    "blockchain": "7C3AED",
    "external": "0EA5E9",
}

TYPE_STYLES: dict[str, dict[str, str]] = {
    "platform": {"fill": "E8F0FF", "line": "3B82F6", "icon": "platform"},
    "database": {"fill": "ECFDF5", "line": "10B981", "icon": "database"},
    "api": {"fill": "EFF6FF", "line": "2563EB", "icon": "api"},
    "user": {"fill": "FFF7ED", "line": "F59E0B", "icon": "users"},
    "external": {"fill": "F0F9FF", "line": "0EA5E9", "icon": "integration"},
    "blockchain": {"fill": "F5F3FF", "line": "7C3AED", "icon": "blockchain"},
}

REGION_POINTS = {
    "americas": {"x": 0.18, "y": 0.40},
    "europe": {"x": 0.47, "y": 0.26},
    "middle east": {"x": 0.58, "y": 0.36},
    "asia-pacific": {"x": 0.77, "y": 0.42},
    "africa": {"x": 0.49, "y": 0.52},
}


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", ""))


def _to_emu(value: Any):
    if isinstance(value, (int, float)):
        if 0 <= float(value) <= 20:
            return Inches(float(value))
        return int(value)
    return value


def _position(position: dict[str, Any] | tuple[Any, Any, Any, Any]):
    if isinstance(position, dict):
        return _to_emu(position["left"]), _to_emu(position["top"]), _to_emu(position["width"]), _to_emu(position["height"])
    left, top, width, height = position
    return _to_emu(left), _to_emu(top), _to_emu(width), _to_emu(height)


def _xy(position: dict[str, Any] | tuple[Any, Any]):
    if isinstance(position, dict):
        return _to_emu(position["left"]), _to_emu(position["top"])
    left, top = position
    return _to_emu(left), _to_emu(top)


def _text_box(slide, left, top, width, height, text: str, size=12, color="1E293B", bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Figtree"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _load_logo_registry() -> dict[str, str]:
    if not LOGO_REGISTRY_PATH.exists():
        return {}
    raw = json.loads(LOGO_REGISTRY_PATH.read_text(encoding="utf-8"))
    return {str(k).strip().lower(): str(v) for k, v in raw.items()}


def _load_icon_registry() -> dict[str, str]:
    if not ICON_REGISTRY_PATH.exists():
        return {}
    return json.loads(ICON_REGISTRY_PATH.read_text(encoding="utf-8"))


def _resolve_asset(path_str: str | None, base_dir: Path | None = None) -> Path | None:
    if not path_str:
        return None
    path = Path(path_str)
    if path.is_absolute() and path.exists():
        return path
    candidates = []
    if base_dir:
        candidates.append((base_dir / path).resolve())
    candidates.extend([
        (ROOT / path).resolve(),
        (ROOT / "ppt-maker" / path).resolve(),
        (ROOT / "ppt-maker" / "assets" / path).resolve(),
        (LOGO_DIR / path.name).resolve(),
    ])
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return None


def _emu_to_px(value: Any) -> int | None:
    if value is None:
        return None
    try:
        # Use 192 DPI (2× standard 96 DPI) so embedded bitmaps are Retina-ready.
        # At 96 DPI, a typical image placeholder embeds at ~564×502 px, which looks
        # soft on HiDPI displays. At 192 DPI the same slot embeds at ~1128×1004 px,
        # matching the physical pixel density of Retina screens.
        return max(1, int(round(int(value) / 4762)))
    except Exception:
        return None


def _parse_background_color(background: str | None) -> tuple[int, int, int, int]:
    if not background or str(background).lower() == 'transparent':
        return (255, 255, 255, 0)
    try:
        rgb = ImageColor.getrgb(str(background))
        if len(rgb) == 4:
            return rgb
        return (*rgb, 255)
    except Exception:
        return (255, 255, 255, 255)


def prepare_image_for_placeholder(
    path: str | Path | None,
    *,
    width_emu=None,
    height_emu=None,
    fit_mode: str = 'contain',
    background: str = 'transparent',
    round_corners: bool = False,
    radius_px: int = 8,
    cache_dir: Path | None = None,
) -> Path | None:
    source = Path(path) if path else None
    if not source or not source.exists():
        return None

    target_width_px = _emu_to_px(width_emu)
    target_height_px = _emu_to_px(height_emu)
    cache_root = cache_dir or (ROOT / 'ppt-maker' / 'output' / '.prepared-image-cache')
    cache_root.mkdir(parents=True, exist_ok=True)
    fingerprint = hashlib.sha1(
        f"{source.resolve()}::{source.stat().st_mtime_ns}::{source.stat().st_size}::{target_width_px}::{target_height_px}::{fit_mode}::{background}::{round_corners}::{radius_px}".encode('utf-8')
    ).hexdigest()[:16]
    output_path = cache_root / f"{source.stem}-{fit_mode}-{fingerprint}.png"
    if output_path.exists():
        return output_path

    with Image.open(source) as image:
        image = image.convert('RGBA')
        output = image
        if target_width_px and target_height_px:
            src_w, src_h = image.size
            if fit_mode == 'cover':
                scale = max(target_width_px / src_w, target_height_px / src_h)
            else:
                scale = min(target_width_px / src_w, target_height_px / src_h)
            resized = image.resize(
                (max(1, int(round(src_w * scale))), max(1, int(round(src_h * scale)))),
                Image.Resampling.LANCZOS,
            )
            if fit_mode == 'cover':
                left = max(0, int(round((resized.width - target_width_px) / 2)))
                top = max(0, int(round((resized.height - target_height_px) / 2)))
                output = resized.crop((left, top, left + target_width_px, top + target_height_px))
            else:
                output = Image.new('RGBA', (target_width_px, target_height_px), _parse_background_color(background))
                paste_left = int(round((target_width_px - resized.width) / 2))
                paste_top = int(round((target_height_px - resized.height) / 2))
                output.alpha_composite(resized, (paste_left, paste_top))

        if round_corners:
            mask = Image.new('L', output.size, 0)
            draw = ImageDraw.Draw(mask)
            max_radius = max(1, min(radius_px, output.size[0] // 2, output.size[1] // 2))
            draw.rounded_rectangle((0, 0, output.size[0] - 1, output.size[1] - 1), radius=max_radius, fill=255)
            output.putalpha(mask)
        output.save(output_path, format='PNG')
    return output_path


def prepare_rounded_image(path: str | Path | None, *, radius_px: int = 8, width_emu=None, height_emu=None, cache_dir: Path | None = None) -> Path | None:
    return prepare_image_for_placeholder(
        path,
        width_emu=width_emu,
        height_emu=height_emu,
        fit_mode='cover',
        background='transparent',
        round_corners=True,
        radius_px=radius_px,
        cache_dir=cache_dir,
    )


def add_picture_fit(slide, path: str | Path | None, left, top, width=None, height=None, *, round_corners: bool = False, radius_px: int = 8, fit_mode: str = 'contain', background: str = 'transparent', cache_dir: Path | None = None):
    """Add an image without aspect-ratio distortion.

    - contain: whole image visible, centered in box
    - cover: fills box, centered, cropped if needed

    When rounded corners or a non-transparent background are requested, we
    pre-compose an aspect-ratio-safe bitmap at the target box size and embed
    that result. Otherwise we place the original asset with explicit fit math
    instead of forcing width+height blindly.
    """
    picture_path = Path(path) if path else None
    if not picture_path or not picture_path.exists():
        return None

    if width is None and height is None:
        return slide.shapes.add_picture(str(picture_path), left, top)

    if width is None or height is None:
        return slide.shapes.add_picture(str(picture_path), left, top, width=width, height=height)

    needs_prepared_bitmap = round_corners or str(background).lower() != 'transparent'
    if needs_prepared_bitmap:
        prepared = prepare_image_for_placeholder(
            picture_path,
            width_emu=width,
            height_emu=height,
            fit_mode=fit_mode,
            background=background,
            round_corners=round_corners,
            radius_px=radius_px,
            cache_dir=cache_dir,
        )
        if prepared:
            return slide.shapes.add_picture(str(prepared), left, top, width=width, height=height)

    with Image.open(picture_path) as image:
        src_w, src_h = image.size

    if src_w <= 0 or src_h <= 0:
        return slide.shapes.add_picture(str(picture_path), left, top, width=width, height=height)

    box_w = int(width)
    box_h = int(height)
    image_ratio = src_w / src_h
    box_ratio = box_w / box_h if box_h else image_ratio

    if fit_mode == 'cover':
        if image_ratio > box_ratio:
            picture = slide.shapes.add_picture(str(picture_path), left, top, height=height)
            visible_fraction = min(1.0, box_w / picture.width) if picture.width else 1.0
            crop = max(0.0, (1.0 - visible_fraction) / 2.0)
            picture.crop_left = crop
            picture.crop_right = crop
            picture.left = left
            picture.top = top
            picture.width = width
            picture.height = height
            return picture
        picture = slide.shapes.add_picture(str(picture_path), left, top, width=width)
        visible_fraction = min(1.0, box_h / picture.height) if picture.height else 1.0
        crop = max(0.0, (1.0 - visible_fraction) / 2.0)
        picture.crop_top = crop
        picture.crop_bottom = crop
        picture.left = left
        picture.top = top
        picture.width = width
        picture.height = height
        return picture

    if image_ratio > box_ratio:
        render_w = box_w
        render_h = max(1, int(round(render_w / image_ratio)))
    else:
        render_h = box_h
        render_w = max(1, int(round(render_h * image_ratio)))

    render_left = int(left + (box_w - render_w) / 2)
    render_top = int(top + (box_h - render_h) / 2)
    return slide.shapes.add_picture(str(picture_path), render_left, render_top, width=render_w, height=render_h)


def _safe_picture(slide, path: Path | None, left, top, width=None, height=None, *, round_corners: bool = False, radius_px: int = 8, fit_mode: str = 'cover', background: str = 'transparent'):
    return add_picture_fit(
        slide,
        path,
        left,
        top,
        width=width,
        height=height,
        round_corners=round_corners,
        radius_px=radius_px,
        fit_mode=fit_mode,
        background=background,
    )


def _icon_path(icon_name: str | None) -> Path | None:
    if not icon_name:
        return None
    registry = _load_icon_registry()
    icon_rel = registry.get(icon_name)
    if not icon_rel:
        return None
    return (ICON_REGISTRY_PATH.parent / Path(icon_rel).name).resolve()


def add_logo(slide, company_name: str, position, max_height=None, registry_path: Path | None = None, base_dir: Path | None = None):
    """Add a company logo scaled to fit inside the target box."""
    left, top, width, height = _position(position)
    registry_file = registry_path or LOGO_REGISTRY_PATH
    registry = _load_logo_registry() if registry_file == LOGO_REGISTRY_PATH else {
        str(k).strip().lower(): str(v) for k, v in json.loads(registry_file.read_text(encoding="utf-8")).items()
    }
    key = company_name.strip().lower()
    logo_path = _resolve_asset(registry.get(key), base_dir=base_dir)
    if not logo_path:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb(BRAND_COLORS["white"])
        box.line.color.rgb = _rgb(BRAND_COLORS["mid_gray"])
        box.line.width = Pt(1)
        _text_box(slide, left + Inches(0.08), top, width - Inches(0.16), height, company_name, size=11, color=BRAND_COLORS["gray"], bold=True, align=PP_ALIGN.CENTER)
        return {"shape": box, "missing": True}

    target_height = min(height, max_height or height)
    picture = slide.shapes.add_picture(str(logo_path), left, top, height=target_height)
    if picture.width > width:
        picture.width = width
    picture.left = left + int((width - picture.width) / 2)
    picture.top = top + int((height - picture.height) / 2)
    return {"picture": picture, "missing": False, "path": str(logo_path)}


def add_logo_grid(slide, company_names: list[str], area, max_per_row: int = 4, base_dir: Path | None = None):
    """Add a trusted-by logo grid."""
    left, top, width, height = _position(area)
    names = company_names or []
    if not names:
        return []
    cols = max(1, min(max_per_row, len(names)))
    rows = math.ceil(len(names) / cols)
    gap_x = Inches(0.18)
    gap_y = Inches(0.14)
    cell_width = (width - gap_x * (cols - 1)) / cols
    cell_height = (height - gap_y * (rows - 1)) / rows
    added = []
    for idx, company in enumerate(names):
        row = idx // cols
        col = idx % cols
        x = left + col * (cell_width + gap_x)
        y = top + row * (cell_height + gap_y)
        added.append(add_logo(slide, company, (x, y, cell_width, cell_height), max_height=cell_height * 0.68, base_dir=base_dir))
    return added


def add_logo_strip(slide, company_names: list[str], y_position, base_dir: Path | None = None):
    """Add a bottom-of-slide horizontal logo strip."""
    if not company_names:
        return []
    y_position = _to_emu(y_position)
    slide_width = slide.part.package.presentation_part.presentation.slide_width
    margin = Inches(0.7)
    gap = Inches(0.18)
    usable_width = slide_width - (2 * margin)
    cell_width = (usable_width - gap * (len(company_names) - 1)) / max(len(company_names), 1)
    cell_height = Inches(0.42)
    added = []
    for idx, company in enumerate(company_names):
        x = margin + idx * (cell_width + gap)
        added.append(add_logo(slide, company, (x, y_position, cell_width, cell_height), max_height=cell_height * 0.8, base_dir=base_dir))
    return added


def _component_position(component: dict[str, Any], area_box, index: int, total: int):
    left, top, width, height = area_box
    hint = (component.get("position_hint") or "").lower()
    if hint in {"left", "west"}:
        return left + width * 0.05, top + height * 0.35
    if hint in {"right", "east"}:
        return left + width * 0.72, top + height * 0.35
    if hint in {"top", "north"}:
        return left + width * 0.38, top + height * 0.05
    if hint in {"bottom", "south"}:
        return left + width * 0.38, top + height * 0.70
    if hint in {"center", "middle"}:
        return left + width * 0.38, top + height * 0.35
    if total <= 1:
        return left + width * 0.38, top + height * 0.35
    cols = min(3, total)
    rows = math.ceil(total / cols)
    col = index % cols
    row = index // cols
    box_w = width * 0.24
    box_h = height * 0.18
    gap_x = (width - cols * box_w) / max(cols + 1, 1)
    gap_y = (height - rows * box_h) / max(rows + 1, 1)
    return left + gap_x + col * (box_w + gap_x), top + gap_y + row * (box_h + gap_y)


def add_architecture_diagram(slide, components: list[dict[str, Any]], connections: list[dict[str, Any]], area):
    """Render an editable architecture diagram using native shapes and connectors."""
    left, top, width, height = _position(area)
    area_box = (left, top, width, height)
    nodes: dict[str, dict[str, Any]] = {}
    rendered = []
    for idx, component in enumerate(components or []):
        comp_type = str(component.get("type", "platform")).lower()
        style = TYPE_STYLES.get(comp_type, TYPE_STYLES["platform"])
        x, y = _component_position(component, area_box, idx, len(components or []))
        box_w = component.get("width") or width * 0.24
        box_h = component.get("height") or height * 0.18
        shape_kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if comp_type == "database":
            shape_kind = MSO_AUTO_SHAPE_TYPE.CAN
        elif comp_type == "user":
            shape_kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        elif comp_type == "blockchain":
            shape_kind = MSO_AUTO_SHAPE_TYPE.HEXAGON
        elif comp_type == "external":
            shape_kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        node = slide.shapes.add_shape(shape_kind, int(x), int(y), int(box_w), int(box_h))
        node.fill.solid()
        node.fill.fore_color.rgb = _rgb(style["fill"])
        node.line.color.rgb = _rgb(style["line"])
        node.line.width = Pt(1.6)
        icon_path = _icon_path(style.get("icon"))
        if icon_path:
            _safe_picture(slide, icon_path, int(x + Inches(0.08)), int(y + Inches(0.08)), width=Inches(0.22), height=Inches(0.22))
        label_left = x + Inches(0.34)
        label_width = max(box_w - Inches(0.42), Inches(0.8))
        title = str(component.get("name", "Component"))
        subtitle = component.get("subtitle")
        title_box = _text_box(slide, int(label_left), int(y + Inches(0.04)), int(label_width), int(box_h * 0.45), title, size=12, color=BRAND_COLORS["dark"], bold=True)
        if subtitle:
            _text_box(slide, int(label_left), int(y + box_h * 0.45), int(label_width), int(box_h * 0.28), str(subtitle), size=9, color=BRAND_COLORS["gray"])
        key = str(component.get("id") or component.get("name") or f"component_{idx}")
        nodes[key] = {
            "shape": node,
            "center": (int(x + box_w / 2), int(y + box_h / 2)),
            "left": int(x),
            "top": int(y),
            "width": int(box_w),
            "height": int(box_h),
            "title_box": title_box,
        }
        rendered.append(node)

    for connection in connections or []:
        source_key = str(connection.get("from") or connection.get("source") or "")
        target_key = str(connection.get("to") or connection.get("target") or "")
        if source_key not in nodes or target_key not in nodes:
            continue
        source = nodes[source_key]
        target = nodes[target_key]
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, source["center"][0], source["center"][1], target["center"][0], target["center"][1])
        connector.line.color.rgb = _rgb(connection.get("color", BRAND_COLORS["blue"]))
        connector.line.width = Pt(connection.get("width", 1.8))
        if hasattr(connector.line, "end_arrowhead"):
            connector.line.end_arrowhead = True
        rendered.append(connector)
        if connection.get("label"):
            label_x = int((source["center"][0] + target["center"][0]) / 2 - Inches(0.45))
            label_y = int((source["center"][1] + target["center"][1]) / 2 - Inches(0.12))
            label = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, label_x, label_y, Inches(0.9), Inches(0.26))
            label.fill.solid()
            label.fill.fore_color.rgb = _rgb(BRAND_COLORS["white"])
            label.line.color.rgb = _rgb(BRAND_COLORS["mid_gray"])
            label.line.width = Pt(0.75)
            _text_box(slide, label_x + Inches(0.04), label_y, Inches(0.82), Inches(0.26), str(connection["label"]), size=8, color=BRAND_COLORS["gray"], align=PP_ALIGN.CENTER)
    return rendered


def add_flow_diagram(slide, steps: list[dict[str, Any] | str], area):
    """Render a linear or branched flow diagram."""
    left, top, width, height = _position(area)
    normalized = []
    for idx, step in enumerate(steps or []):
        if isinstance(step, str):
            normalized.append({"id": f"step_{idx}", "title": step})
        else:
            item = dict(step)
            item.setdefault("id", f"step_{idx}")
            normalized.append(item)
    if not normalized:
        return []

    root_steps = [s for s in normalized if not s.get("branch_of")]
    step_width = min(Inches(1.8), width / max(len(root_steps), 1) - Inches(0.15))
    step_height = Inches(0.72)
    top_y = top + height * 0.15
    nodes = {}
    rendered = []
    for idx, step in enumerate(root_steps):
        x = left + idx * (step_width + Inches(0.22))
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, int(x), int(top_y), int(step_width), int(step_height))
        node.fill.solid()
        node.fill.fore_color.rgb = _rgb("EFF6FF")
        node.line.color.rgb = _rgb(BRAND_COLORS["blue"])
        node.line.width = Pt(1.4)
        _text_box(slide, int(x + Inches(0.08)), int(top_y + Inches(0.08)), int(step_width - Inches(0.16)), int(Inches(0.26)), str(step.get("title", "Step")), size=11, color=BRAND_COLORS["dark"], bold=True, align=PP_ALIGN.CENTER)
        if step.get("subtitle"):
            _text_box(slide, int(x + Inches(0.08)), int(top_y + Inches(0.34)), int(step_width - Inches(0.16)), int(Inches(0.18)), str(step.get("subtitle")), size=8, color=BRAND_COLORS["gray"], align=PP_ALIGN.CENTER)
        nodes[step["id"]] = {"center": (int(x + step_width / 2), int(top_y + step_height / 2)), "x": int(x), "y": int(top_y)}
        rendered.append(node)

    for idx in range(len(root_steps) - 1):
        a = nodes[root_steps[idx]["id"]]["center"]
        b = nodes[root_steps[idx + 1]["id"]]["center"]
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, a[0], a[1], b[0], b[1])
        connector.line.color.rgb = _rgb(BRAND_COLORS["blue"])
        connector.line.width = Pt(1.75)
        if hasattr(connector.line, "end_arrowhead"):
            connector.line.end_arrowhead = True
        rendered.append(connector)

    branch_y = top + height * 0.62
    branches = [s for s in normalized if s.get("branch_of")]
    for idx, step in enumerate(branches):
        parent = nodes.get(step.get("branch_of"))
        if not parent:
            continue
        x = parent["x"] + (idx % 2) * (step_width + Inches(0.22)) - Inches(0.6)
        y = branch_y + (idx // 2) * Inches(0.9)
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, int(x), int(y), int(step_width), int(step_height))
        node.fill.solid()
        node.fill.fore_color.rgb = _rgb("FFF7ED" if str(step.get("branch_label", "")).lower() == "no" else "ECFDF5")
        node.line.color.rgb = _rgb(BRAND_COLORS["amber"] if str(step.get("branch_label", "")).lower() == "no" else BRAND_COLORS["green"])
        node.line.width = Pt(1.3)
        title = step.get("title", "Branch")
        if step.get("branch_label"):
            title = f"{step['branch_label']}: {title}"
        _text_box(slide, int(x + Inches(0.08)), int(y + Inches(0.08)), int(step_width - Inches(0.16)), int(step_height - Inches(0.16)), str(title), size=10, color=BRAND_COLORS["dark"], bold=True, align=PP_ALIGN.CENTER)
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, parent["center"][0], parent["center"][1], int(x + step_width / 2), int(y))
        connector.line.color.rgb = _rgb(BRAND_COLORS["gray"])
        connector.line.width = Pt(1.4)
        if hasattr(connector.line, "end_arrowhead"):
            connector.line.end_arrowhead = True
        rendered.extend([node, connector])
    return rendered


def add_laptop_frame(slide, screenshot_path: str, position, width, base_dir: Path | None = None):
    """Add a clean laptop device frame around a screenshot."""
    left, top = _xy(position)
    width = _to_emu(width)
    frame_height = int(width * 0.68)
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, int(frame_height))
    body.fill.solid()
    body.fill.fore_color.rgb = _rgb("334155")
    body.line.fill.background()
    bezel = Inches(0.12)
    screen_left = left + bezel
    screen_top = top + bezel
    screen_width = width - 2 * bezel
    screen_height = int(frame_height - 2 * bezel - Inches(0.14))
    screen = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, int(screen_left), int(screen_top), int(screen_width), int(screen_height))
    screen.fill.solid()
    screen.fill.fore_color.rgb = _rgb("0F172A")
    screen.line.fill.background()
    screenshot = _safe_picture(slide, _resolve_asset(screenshot_path, base_dir=base_dir), int(screen_left + Inches(0.04)), int(screen_top + Inches(0.04)), width=int(screen_width - Inches(0.08)), height=int(screen_height - Inches(0.08)), round_corners=True, fit_mode='contain', background='white')
    base = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, left + int(width * 0.2), int(top + frame_height - Inches(0.02)), int(width * 0.6), int(Inches(0.20)))
    base.rotation = 180
    base.fill.solid()
    base.fill.fore_color.rgb = _rgb("94A3B8")
    base.line.fill.background()
    return {"body": body, "screen": screen, "picture": screenshot, "base": base}


def add_phone_frame(slide, screenshot_path: str, position, height, base_dir: Path | None = None):
    """Add a phone frame with rounded corners and a notch."""
    left, top = _xy(position)
    height = _to_emu(height)
    width = int(height * 0.49)
    phone = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    phone.fill.solid()
    phone.fill.fore_color.rgb = _rgb("111827")
    phone.line.fill.background()
    inset = Inches(0.08)
    screen_left = left + inset
    screen_top = top + inset + Inches(0.08)
    screen_width = width - 2 * inset
    screen_height = height - 2 * inset - Inches(0.14)
    screen = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, int(screen_left), int(screen_top), int(screen_width), int(screen_height))
    screen.fill.solid()
    screen.fill.fore_color.rgb = _rgb("F8FAFC")
    screen.line.fill.background()
    notch = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, int(left + width * 0.30), int(top + Inches(0.05)), int(width * 0.40), int(Inches(0.10)))
    notch.fill.solid()
    notch.fill.fore_color.rgb = _rgb("0F172A")
    notch.line.fill.background()
    screenshot = _safe_picture(slide, _resolve_asset(screenshot_path, base_dir=base_dir), int(screen_left), int(screen_top), width=int(screen_width), height=int(screen_height), round_corners=True, fit_mode='contain', background='white')
    return {"phone": phone, "screen": screen, "picture": screenshot, "notch": notch}


def add_browser_frame(slide, screenshot_path: str, position, width, base_dir: Path | None = None):
    """Add a browser chrome frame with screenshot content below."""
    left, top = _xy(position)
    width = _to_emu(width)
    total_height = int(width * 0.62)
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, total_height)
    outer.fill.solid()
    outer.fill.fore_color.rgb = _rgb(BRAND_COLORS["white"])
    outer.line.color.rgb = _rgb(BRAND_COLORS["mid_gray"])
    outer.line.width = Pt(1)
    chrome_h = Inches(0.34)
    chrome = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, chrome_h)
    chrome.fill.solid()
    chrome.fill.fore_color.rgb = _rgb("E2E8F0")
    chrome.line.fill.background()
    dot_colors = ["EF4444", "F59E0B", "10B981"]
    for idx, color in enumerate(dot_colors):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, int(left + Inches(0.10) + idx * Inches(0.12)), int(top + Inches(0.10)), int(Inches(0.06)), int(Inches(0.06)))
        dot.fill.solid()
        dot.fill.fore_color.rgb = _rgb(color)
        dot.line.fill.background()
    url_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, int(left + Inches(0.55)), int(top + Inches(0.07)), int(width - Inches(0.75)), int(Inches(0.12)))
    url_bar.fill.solid()
    url_bar.fill.fore_color.rgb = _rgb(BRAND_COLORS["white"])
    url_bar.line.color.rgb = _rgb("CBD5E1")
    url_bar.line.width = Pt(0.7)
    _safe_picture(slide, _resolve_asset(screenshot_path, base_dir=base_dir), int(left + Inches(0.03)), int(top + chrome_h + Inches(0.02)), width=int(width - Inches(0.06)), height=int(total_height - chrome_h - Inches(0.05)), round_corners=True, fit_mode='contain', background='white')
    return {"outer": outer, "chrome": chrome, "url_bar": url_bar}


def add_world_map_markers(slide, locations: list[dict[str, Any] | str], area):
    """Add a schematic world map with regional markers and labels."""
    left, top, width, height = _position(area)
    background = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    background.fill.solid()
    background.fill.fore_color.rgb = _rgb("F8FAFC")
    background.line.color.rgb = _rgb("E2E8F0")
    background.line.width = Pt(1)

    regions = [
        ("Americas", 0.06, 0.22, 0.22, 0.48),
        ("Europe", 0.38, 0.12, 0.16, 0.22),
        ("Middle East", 0.54, 0.24, 0.12, 0.18),
        ("Asia-Pacific", 0.67, 0.18, 0.24, 0.40),
        ("Africa", 0.42, 0.34, 0.14, 0.24),
    ]
    for label, rx, ry, rw, rh in regions:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CLOUD, int(left + width * rx), int(top + height * ry), int(width * rw), int(height * rh))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb("EAF2FF")
        shape.line.color.rgb = _rgb("D2E3FC")
        shape.line.width = Pt(0.8)
        _text_box(slide, int(left + width * rx), int(top + height * (ry + rh) - Inches(0.10)), int(width * rw), int(Inches(0.18)), label, size=8, color=BRAND_COLORS["gray"], align=PP_ALIGN.CENTER)

    rendered = [background]
    offsets: dict[str, int] = {}
    for item in locations or []:
        if isinstance(item, str):
            item = {"label": item}
        region = str(item.get("region") or item.get("group") or item.get("label") or "Europe").lower()
        point = REGION_POINTS.get(region, REGION_POINTS["europe"])
        idx = offsets.get(region, 0)
        offsets[region] = idx + 1
        jitter_x = (idx % 3) * 0.018
        jitter_y = (idx // 3) * 0.04
        x = left + width * (point["x"] + jitter_x)
        y = top + height * (point["y"] + jitter_y)
        marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, int(x), int(y), int(Inches(0.12)), int(Inches(0.12)))
        marker.fill.solid()
        marker.fill.fore_color.rgb = _rgb(item.get("color", BRAND_COLORS["blue"]))
        marker.line.fill.background()
        label = item.get("label") or item.get("name") or region.title()
        _text_box(slide, int(x + Inches(0.12)), int(y - Inches(0.02)), int(Inches(1.0)), int(Inches(0.18)), str(label), size=8, color=BRAND_COLORS["dark"])
        rendered.append(marker)
    return rendered


def _initials(name: str) -> str:
    bits = [bit[0].upper() for bit in name.split() if bit]
    return "".join(bits[:2]) or "?"


def add_team_grid(slide, team_members: list[dict[str, Any]], area, base_dir: Path | None = None):
    """Add a team grid with photos or initials placeholders."""
    left, top, width, height = _position(area)
    people = team_members or []
    if not people:
        return []
    cols = min(4, max(2, len(people))) if len(people) > 1 else 1
    rows = math.ceil(len(people) / cols)
    gap_x = Inches(0.18)
    gap_y = Inches(0.24)
    card_width = (width - gap_x * (cols - 1)) / cols
    card_height = (height - gap_y * (rows - 1)) / rows
    photo_size = min(card_width * 0.68, card_height * 0.48)
    rendered = []
    for idx, member in enumerate(people):
        row = idx // cols
        col = idx % cols
        x = left + col * (card_width + gap_x)
        y = top + row * (card_height + gap_y)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, int(x), int(y), int(card_width), int(card_height))
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(BRAND_COLORS["white"])
        card.line.color.rgb = _rgb("E2E8F0")
        card.line.width = Pt(1)
        photo_x = x + (card_width - photo_size) / 2
        photo_y = y + Inches(0.12)
        photo_path = _resolve_asset(member.get("photo_path"), base_dir=base_dir)
        if photo_path:
            _safe_picture(slide, photo_path, int(photo_x), int(photo_y), width=int(photo_size), height=int(photo_size), round_corners=True)
            ring = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, int(photo_x), int(photo_y), int(photo_size), int(photo_size))
            ring.fill.background()
            ring.line.color.rgb = _rgb("CBD5E1")
            ring.line.width = Pt(1.1)
            rendered.append(ring)
        else:
            avatar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, int(photo_x), int(photo_y), int(photo_size), int(photo_size))
            avatar.fill.solid()
            avatar.fill.fore_color.rgb = _rgb("DBEAFE")
            avatar.line.fill.background()
            _text_box(slide, int(photo_x), int(photo_y + photo_size * 0.22), int(photo_size), int(photo_size * 0.28), _initials(member.get("name", "?")), size=20, color=BRAND_COLORS["blue"], bold=True, align=PP_ALIGN.CENTER)
            rendered.append(avatar)
        _text_box(slide, int(x + Inches(0.08)), int(y + photo_size + Inches(0.22)), int(card_width - Inches(0.16)), int(Inches(0.22)), str(member.get("name", "Name")), size=11, color=BRAND_COLORS["dark"], bold=True, align=PP_ALIGN.CENTER)
        _text_box(slide, int(x + Inches(0.08)), int(y + photo_size + Inches(0.46)), int(card_width - Inches(0.16)), int(Inches(0.30)), str(member.get("title", "")), size=9, color=BRAND_COLORS["gray"], align=PP_ALIGN.CENTER)
        rendered.append(card)
    return rendered


def add_photo_with_caption(slide, photo_path: str, caption: str, position, base_dir: Path | None = None):
    """Add a photo with a caption bar."""
    left, top, width, height = _position(position)
    caption_height = Inches(0.34)
    picture = _safe_picture(slide, _resolve_asset(photo_path, base_dir=base_dir), left, top, width=width, height=height - caption_height, round_corners=True)
    caption_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top + height - caption_height, width, caption_height)
    caption_bar.fill.solid()
    caption_bar.fill.fore_color.rgb = _rgb("E2E8F0")
    caption_bar.line.fill.background()
    _text_box(slide, int(left + Inches(0.10)), int(top + height - caption_height), int(width - Inches(0.20)), int(caption_height), caption, size=10, color=BRAND_COLORS["dark"], bold=True, align=PP_ALIGN.CENTER)
    return {"picture": picture, "caption_bar": caption_bar}


def add_visual_element(slide, element: dict[str, Any], base_dir: Path | None = None):
    """Dispatch JSON-configured visual elements from build-pptx.py."""
    element_type = str(element.get("type", "")).lower()
    if element_type == "logo":
        return add_logo(slide, element["company"], element["position"], element.get("max_height"), base_dir=base_dir)
    if element_type == "logo_grid":
        return add_logo_grid(slide, element.get("companies", []), element["area"], element.get("max_per_row", 4), base_dir=base_dir)
    if element_type == "logo_strip":
        return add_logo_strip(slide, element.get("companies", []), element["y_position"], base_dir=base_dir)
    if element_type == "architecture":
        return add_architecture_diagram(slide, element.get("components", []), element.get("connections", []), element["area"])
    if element_type == "flow":
        return add_flow_diagram(slide, element.get("steps", []), element["area"])
    if element_type == "device_mockup":
        device = str(element.get("device", "browser")).lower()
        position = element.get("position") or (element.get("left"), element.get("top"))
        if device == "laptop":
            return add_laptop_frame(slide, element["screenshot"], position, element["width"], base_dir=base_dir)
        if device == "phone":
            return add_phone_frame(slide, element["screenshot"], position, element["height"], base_dir=base_dir)
        return add_browser_frame(slide, element["screenshot"], position, element["width"], base_dir=base_dir)
    if element_type == "world_map":
        return add_world_map_markers(slide, element.get("locations", []), element["area"])
    if element_type == "team_grid":
        return add_team_grid(slide, element.get("team_members", []), element["area"], base_dir=base_dir)
    if element_type == "photo_caption":
        return add_photo_with_caption(slide, element["photo_path"], element.get("caption", ""), element["position"], base_dir=base_dir)
    raise KeyError(f"Unsupported visual element type: {element_type}")
