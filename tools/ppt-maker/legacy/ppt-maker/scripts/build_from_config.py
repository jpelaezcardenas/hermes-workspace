#!/usr/bin/env python3
from __future__ import annotations
"""
Build PPTX from JSON config using slide bank + Master Template.

Workflow:
1. Read config JSON defining slides + content
2. For each slide: copy bank file uniquely, edit content (no styling)
3. Extract only slide XML from edited files
4. Merge into Master Template (delete original slides, add ours)
"""

import json
import math
import re
import subprocess
import zipfile
import os
import shutil
import tempfile
import sys
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.xmlchemy import OxmlElement

from image_helpers import add_picture_fit, prepare_image_for_placeholder, prepare_rounded_image

ROOT = Path(__file__).resolve().parent.parent
FIGTREE_FONT = ROOT.parent.parent / "assets" / "fonts" / "figtree" / "ttf" / "Figtree-SemiBold.ttf"
FIGTREE_BOLD_FONT = ROOT.parent.parent / "assets" / "fonts" / "figtree" / "ttf" / "Figtree-Bold.ttf"
BULLET_PREFIXES = ('•', '-', '*', '–', '—')
MAIN_TITLE_WORD_LIMIT = 5
COLUMN_TITLE_WORD_LIMIT = 4
COLUMN_TITLE_FIELD_SPECS = {}
EMU_PER_INCH = 914400
BOLD_MARKER_RE = re.compile(r"(\*\*[^*].+?\*\*)")
# Single-asterisk emphasis pattern (*text*) — strip these to plain text since
# PowerPoint body text doesn't use italics in SettleMint brand styling.
# Must NOT match ** (bold) or bullet-prefix * at line start.
SINGLE_EMPHASIS_RE = re.compile(r'(?<!\*)\*(?!\*)([^*\n]+?)(?<!\*)\*(?!\*)')
FIT_REPORT_VERSION = "2026-03-20"
APPENDIX_MIN_DIAGRAMS = 4
APPENDIX_MAX_DIAGRAMS = 10
AGENDA_TEXT_PLACEHOLDER_START = 11
AGENDA_ICON_PLACEHOLDER_START = 18
AGENDA_ITEM_COUNT = 7
AGENDA_SHORT_PHRASE_WORD_LIMIT = 6

PLACEHOLDER_CLASS_PROFILES = {
    'screenshot-browser-frame': {
        'ratio': 1.955,
        'render_width': 2800,
        'slot_px': (1661, 849),   # width × height at 192 DPI
        'legacy_aliases': {'browser_landscape', 'screenshot', 'browser-frame'},
    },
    'equal': {
        'ratio': 1.125,
        'render_width': 1800,
        'slot_px': (1129, 1004),
        'legacy_aliases': {'square', 'equal'},
    },
    'wide': {
        'ratio': 1.678,
        'render_width': 2400,
        'slot_px': (1684, 1004),
        'legacy_aliases': {'landscape', 'wide'},
    },
    'one-third': {
        'ratio': 0.670,
        'render_width': 1400,
        'slot_px': (672, 1004),
        'legacy_aliases': {'portrait', 'one-third', 'one_third'},
    },
    'full-width': {
        'ratio': 2.866,
        'render_width': 3200,
        'slot_px': (2368, 826),
        'legacy_aliases': {'panorama', 'full-width', 'full_width'},
    },
}

BODY_FIT_PROFILES = {
    'body_default': {
        'font_pt': 12, 'char_width_em': 0.51, 'line_height': 1.22,
        'usable_width': 0.88, 'usable_height': 0.84,
        'warn_under': 0.34, 'fail_under': 0.18,
        'warn_over': 1.0, 'fail_over': 1.18,
        'target_band': [0.48, 0.90],
        'max_size': 18,
    },
    'body_dense': {
        'font_pt': 11.5, 'char_width_em': 0.50, 'line_height': 1.20,
        'usable_width': 0.90, 'usable_height': 0.86,
        'warn_under': 0.32, 'fail_under': 0.18,
        'warn_over': 1.03, 'fail_over': 1.18,
        'target_band': [0.46, 0.88],
        'max_size': 18,
    },
    'three_col': {
        'font_pt': 11.5, 'char_width_em': 0.50, 'line_height': 1.18,
        'usable_width': 0.91, 'usable_height': 0.86,
        'warn_under': 0.34, 'fail_under': 0.22,
        'warn_over': 1.0, 'fail_over': 1.14,
        'target_band': [0.50, 0.90],
        'max_size': 18,
    },
    'grid_cell': {
        'font_pt': 11.5, 'char_width_em': 0.50, 'line_height': 1.16,
        'usable_width': 0.91, 'usable_height': 0.86,
        'warn_under': 0.34, 'fail_under': 0.22,
        'warn_over': 1.0, 'fail_over': 1.12,
        'target_band': [0.50, 0.92],
        'max_size': 18,
    },
    'grid_cell_compact': {
        'font_pt': 10.8, 'char_width_em': 0.50, 'line_height': 1.14,
        'usable_width': 0.92, 'usable_height': 0.87,
        'warn_under': 0.30, 'fail_under': 0.18,
        'warn_over': 1.0, 'fail_over': 1.10,
        'target_band': [0.46, 0.92],
        'max_size': 16,
    },
    'screenshot_text': {
        'font_pt': 11.3, 'char_width_em': 0.50, 'line_height': 1.20,
        'usable_width': 0.90, 'usable_height': 0.85,
        'warn_under': 0.34, 'fail_under': 0.20,
        'warn_over': 1.0, 'fail_over': 1.15,
        'target_band': [0.48, 0.90],
        'max_size': 18,
    },
    'picture_narrow_text': {
        'font_pt': 11.2, 'char_width_em': 0.50, 'line_height': 1.20,
        'usable_width': 0.90, 'usable_height': 0.85,
        'warn_under': 0.32, 'fail_under': 0.20,
        'warn_over': 1.0, 'fail_over': 1.16,
        'target_band': [0.46, 0.90],
        'max_size': 18,
    },
}

SLIDE_BANK_REGISTRY_PATH = ROOT / 'slide-bank' / 'REGISTRY.json'
IMAGE_POLICY_CACHE = None


def _load_slide_bank_registry():
    global IMAGE_POLICY_CACHE
    if IMAGE_POLICY_CACHE is None:
        if SLIDE_BANK_REGISTRY_PATH.exists():
            IMAGE_POLICY_CACHE = json.loads(SLIDE_BANK_REGISTRY_PATH.read_text(encoding='utf-8'))
        else:
            IMAGE_POLICY_CACHE = {'slides': []}
    return IMAGE_POLICY_CACHE


def _get_slide_bank_info(bank_file):
    filename = Path(bank_file or '').name
    registry = _load_slide_bank_registry()
    for slide in registry.get('slides', []):
        if slide.get('filename') == filename:
            return slide
    info_path = ROOT / 'slide-bank' / filename.replace('.pptx', '.info.json')
    if info_path.exists():
        return json.loads(info_path.read_text(encoding='utf-8'))
    return {}


def _canonical_placeholder_class(raw_class, bank_file='', placeholder_kind=''):
    raw = str(raw_class or '').strip().lower()
    if raw in PLACEHOLDER_CLASS_PROFILES:
        return raw

    for canonical, profile in PLACEHOLDER_CLASS_PROFILES.items():
        if raw in profile.get('legacy_aliases', set()):
            return canonical

    bank_key = (bank_file or '').lower()
    kind_key = (placeholder_kind or '').lower()
    if 'screenshot' in bank_key or 'browser' in bank_key or 'browser' in kind_key:
        return 'screenshot-browser-frame'
    if 'full-width' in bank_key or 'full visual' in bank_key:
        return 'full-width'
    if 'one-third' in bank_key or 'one third' in bank_key:
        return 'one-third'
    if 'equal' in bank_key:
        return 'equal'
    if 'wide' in bank_key:
        return 'wide'
    return 'wide'


def _closest_placeholder_class_by_ratio(ratio, placeholder_kind=''):
    if 'browser' in str(placeholder_kind or '').lower():
        return 'screenshot-browser-frame'
    candidates = [
        name for name in PLACEHOLDER_CLASS_PROFILES
        if name != 'screenshot-browser-frame'
    ]
    return min(
        candidates,
        key=lambda name: abs(PLACEHOLDER_CLASS_PROFILES[name]['ratio'] - ratio),
    )


def _placeholder_class_for_slot(slot, bank_file, width=None, height=None):
    placeholder_kind = (slot or {}).get('placeholder_kind', '')
    declared = _canonical_placeholder_class(
        (slot or {}).get('aspect_ratio_class'),
        bank_file=bank_file,
        placeholder_kind=placeholder_kind,
    )
    if width and height:
        measured = _closest_placeholder_class_by_ratio(width / height, placeholder_kind=placeholder_kind)
        if declared == 'wide' and measured != declared:
            return measured
        if declared == 'screenshot-browser-frame':
            return declared
        return measured
    return declared


def _placeholder_profile_for_slot(slot, bank_file, width=None, height=None):
    placeholder_class = _placeholder_class_for_slot(slot, bank_file, width=width, height=height)
    return placeholder_class, PLACEHOLDER_CLASS_PROFILES[placeholder_class]


def _slot_specific_diagram_candidates(diagram_source: Path, placeholder_class: str):
    suffix = diagram_source.suffix
    stem = diagram_source.stem
    variants = [
        diagram_source.with_name(f'{stem}--{placeholder_class}{suffix}'),
        diagram_source.with_name(f'{stem}.{placeholder_class}{suffix}'),
        diagram_source.with_name(f'{stem}-{placeholder_class}{suffix}'),
        diagram_source.with_name(f'{stem}_{placeholder_class}{suffix}'),
    ]
    return variants


def _select_diagram_source_for_slot(diagram_source: Path, placeholder_class: str):
    for candidate in _slot_specific_diagram_candidates(diagram_source, placeholder_class):
        if candidate.exists():
            return candidate
    return diagram_source


def _image_slots_for_slide(bank_file):
    info = _get_slide_bank_info(bank_file)
    slots = info.get('image_slot_policy') or []
    if slots:
        normalized_slots = []
        for slot in slots:
            normalized = dict(slot)
            normalized['aspect_ratio_class'] = _canonical_placeholder_class(
                slot.get('aspect_ratio_class'),
                bank_file=bank_file,
                placeholder_kind=slot.get('placeholder_kind', ''),
            )
            normalized_slots.append(normalized)
        return normalized_slots
    for placeholder in info.get('placeholders', []):
        placeholder_type = str(placeholder.get('type', '')).upper()
        if placeholder_type in {'PICTURE', 'CLIPART'}:
            return [{
                'placeholder_idx': placeholder.get('idx'),
                'placeholder_kind': 'generic_image',
                'allowed_asset_kinds': ['generic_image', 'diagram', 'dalp_screenshot'],
                'preferred_asset_kinds': ['generic_image'],
                'fit_mode_by_asset_kind': {
                    'generic_image': 'cover',
                    'diagram': 'cover',
                    'dalp_screenshot': 'cover',
                },
                'background_by_asset_kind': {
                    'generic_image': 'transparent',
                    'diagram': 'transparent',
                    'dalp_screenshot': 'transparent',
                },
                'round_corners': True,
                'aspect_ratio_class': 'wide',
            }]
    return []


def _resolve_asset_path(image_path):
    if not image_path:
        return None
    candidate_paths = [Path(image_path)]
    if not os.path.isabs(str(image_path)):
        candidate_paths.append(ROOT / image_path)
    for candidate in candidate_paths:
        resolved = candidate.resolve()
        if resolved.exists():
            return resolved
    return None


def _asset_kind_from_path(resolved_path: Path | None):
    if not resolved_path:
        return None
    normalized = str(resolved_path).lower()
    if 'templates/dalp-screenshots/' in normalized or '/dalp-screenshots/' in normalized:
        return 'dalp_screenshot'
    if resolved_path.suffix.lower() in {'.mmd', '.mermaid', '.puml', '.plantuml', '.uml'}:
        return 'diagram_source'
    if 'output/diagrams/' in normalized or '/diagrams/' in normalized:
        return 'diagram'
    diagram_sidecars = [
        resolved_path.with_suffix('.mmd'),
        resolved_path.with_suffix('.mermaid'),
        resolved_path.with_suffix('.puml'),
        resolved_path.with_suffix('.plantuml'),
        resolved_path.with_suffix('.uml'),
    ]
    if resolved_path.suffix.lower() in {'.png', '.svg'} and any(sidecar.exists() for sidecar in diagram_sidecars):
        return 'diagram'
    return 'generic_image'


# ============================================================================
# PARALLEL MERMAID PRE-RENDERING CACHE
# ============================================================================
# When parallel rendering is used, this dict maps diagram source basenames to
# their pre-rendered PNG/SVG paths. Populated by _prerender_diagrams_parallel()
# and consumed by _render_diagram_source().
_PARALLEL_RENDER_CACHE: dict[str, dict[str, str | None]] = {}


def _collect_mermaid_sources_from_config(slides: list[dict]) -> list[tuple[Path, str, str | None]]:
    """Scan config slides and collect all Mermaid diagram source files.

    Returns list of (resolved_source_path, variant_name, placeholder_class).
    """
    collected = []
    for slide_cfg in slides:
        image_path = slide_cfg.get('image_path', '')
        if not image_path:
            continue
        resolved = _resolve_asset_path(image_path)
        if not resolved:
            continue
        if resolved.suffix.lower() not in {'.mmd', '.mermaid'}:
            continue
        bank_file = slide_cfg.get('bank_file', '')
        diagram_target_class = slide_cfg.get('diagram_target_class')
        image_slots = _image_slots_for_slide(bank_file)
        slot = image_slots[0] if image_slots else None
        slot_class = None
        if slot:
            slot_class, _ = _placeholder_profile_for_slot(slot, bank_file)
        requested_class = _canonical_placeholder_class(diagram_target_class, bank_file=bank_file) if diagram_target_class else None
        target_class = requested_class or slot_class or 'wide'
        source_for_slot = _select_diagram_source_for_slot(resolved, target_class)
        variant_name = f"{source_for_slot.stem}-{target_class}"
        collected.append((source_for_slot, variant_name, target_class))
    return collected


def _prerender_diagrams_parallel(slides: list[dict]) -> bool:
    """Pre-render all Mermaid diagrams using Python concurrent.futures.

    Spawns up to 4 concurrent mmdc processes (matching the parallel shell script's
    concurrency) but using Python's ThreadPoolExecutor for reliable cross-platform
    parallelism. The shell script's export -f + background subshells have issues
    with array environment variables (CONFIG_ARGS) on macOS.

    Populates _PARALLEL_RENDER_CACHE with {variant_name: {'png': path, 'svg': path}}.
    Returns True if parallel rendering succeeded, False if fallback is needed.
    """
    global _PARALLEL_RENDER_CACHE
    _PARALLEL_RENDER_CACHE = {}

    sources = _collect_mermaid_sources_from_config(slides)
    if not sources:
        return True  # Nothing to render

    from concurrent.futures import ThreadPoolExecutor, as_completed
    from diagrams.engine import DiagramEngine

    output_dir = ROOT / 'output' / 'diagrams'
    output_dir.mkdir(parents=True, exist_ok=True)

    # Deduplicate by variant_name
    variant_map = {}  # variant_name -> (source_path, target_class)
    for source_path, variant_name, target_class in sources:
        if variant_name not in variant_map:
            variant_map[variant_name] = (source_path, target_class)

    print(f"  [parallel] Rendering {len(variant_map)} diagrams with up to 4 concurrent workers...")

    def _render_one(variant_name, source_path, target_class):
        """Render a single diagram using DiagramEngine (thread-safe: each creates its own mmdc process)."""
        try:
            target_profile = PLACEHOLDER_CLASS_PROFILES.get(target_class, PLACEHOLDER_CLASS_PROFILES['wide'])
            engine = DiagramEngine(output_dir=output_dir)
            render_kwargs = {
                'output_format': 'both',
                'width': target_profile.get('render_width', 2400),
            }
            slot_px = target_profile.get('slot_px')
            if slot_px and len(slot_px) == 2:
                render_kwargs['slot_width'] = slot_px[0]
                render_kwargs['slot_height'] = slot_px[1]
            rendered = engine.render_source_file(source_path, filename=variant_name, **render_kwargs)
            if rendered.get('error'):
                return variant_name, False, rendered['error']
            return variant_name, True, rendered
        except Exception as e:
            return variant_name, False, str(e)

    try:
        max_workers = min(4, len(variant_map))
        success_count = 0

        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {
                executor.submit(_render_one, vname, src, tcls): vname
                for vname, (src, tcls) in variant_map.items()
            }
            for future in as_completed(futures):
                vname = futures[future]
                variant_name_result, ok, result_data = future.result()
                if ok:
                    _PARALLEL_RENDER_CACHE[variant_name_result] = {
                        'png': result_data.get('png'),
                        'svg': result_data.get('svg'),
                    }
                    success_count += 1
                else:
                    print(f"  [parallel] WARN: {variant_name_result} failed: {result_data}")

        print(f"  [parallel] {success_count}/{len(variant_map)} diagrams rendered successfully")

        if success_count < len(variant_map):
            print(f"  [parallel] {len(variant_map) - success_count} will fall back to sequential")

        return success_count > 0

    except Exception as e:
        print(f"  [parallel] Error: {e}, falling back to sequential rendering")
        return False


def _render_diagram_source(diagram_source: Path, variant_name: str, render_width: int | None = None,
                            slot_px: tuple | None = None):
    # Check parallel render cache first
    if variant_name in _PARALLEL_RENDER_CACHE:
        cached = _PARALLEL_RENDER_CACHE[variant_name]
        cached_path = cached.get('svg') or cached.get('png')
        if cached_path and Path(cached_path).exists():
            return Path(cached_path)

    from diagrams.engine import DiagramEngine

    # Always render a fresh diagram for the current presentation context.
    # Do not reuse previously rendered files from output/diagrams/, even if a
    # matching filename already exists from an older deck.
    engine = DiagramEngine(output_dir=ROOT / 'output' / 'diagrams')
    render_kwargs = {'output_format': 'both'}
    if render_width:
        render_kwargs['width'] = render_width
    if slot_px and len(slot_px) == 2:
        render_kwargs['slot_width']  = slot_px[0]
        render_kwargs['slot_height'] = slot_px[1]
    rendered = engine.render_source_file(diagram_source, filename=variant_name, **render_kwargs)
    if rendered.get('error'):
        fallback_svg = diagram_source.with_suffix('.svg')
        fallback_png = diagram_source.with_suffix('.png')
        target_svg = ROOT / 'output' / 'diagrams' / f'{variant_name}.svg'
        target_png = ROOT / 'output' / 'diagrams' / f'{variant_name}.png'
        if fallback_svg.exists():
            shutil.copy2(fallback_svg, target_svg)
        if fallback_png.exists():
            shutil.copy2(fallback_png, target_png)
        if target_svg.exists() or target_png.exists():
            return target_svg if target_svg.exists() else target_png
        raise RuntimeError(f"Diagram render failed for {diagram_source.name}: {rendered['error']}")
    rendered_path = rendered.get('svg') or rendered.get('png')
    return Path(rendered_path) if rendered_path else None


def _is_thank_you_slide(bank_file):
    return 'slide-024-thank-you' in (bank_file or '').lower()


def _is_appendix_section_slide(bank_file, content):
    if 'slide-003' not in (bank_file or '').lower():
        return False
    title = str((content or {}).get('title') or '').strip().lower()
    subtitle = str((content or {}).get('subtitle') or '').strip().lower()
    return title == 'appendix' or subtitle == 'appendix'


def _diagram_asset_kind(image_path):
    resolved = _resolve_asset_path(image_path)
    asset_kind = _asset_kind_from_path(resolved) if resolved else None
    return asset_kind if asset_kind in {'diagram', 'diagram_source'} else None


def _validate_table_text_requirements(bank_file, content):
    # Enforce non-empty body text for ALL table slides (slide-006 and slide-007 both have a BODY_11 text box above the table)
    _table_slide_prefixes = ('slide-006', 'slide-007')
    bank_lower = (bank_file or '').lower()
    if not any(p in bank_lower for p in _table_slide_prefixes):
        return
    if not content.get('table_data'):
        return
    body = content.get('body')
    if not isinstance(body, str) or not body.strip():
        raise ValueError(
            f"{bank_file} requires a non-empty 'body' field when table_data is present; "
            f"the text box above the table must always be filled with a brief description or context sentence."
        )


def _validate_mermaid_workflow(slides):
    thank_you_index = next((i for i, slide in enumerate(slides) if _is_thank_you_slide(slide.get('bank_file', ''))), None)
    if thank_you_index is None:
        return

    appendix_diagram_indices = []
    for idx, slide in enumerate(slides):
        asset_kind = _diagram_asset_kind(slide.get('image_path'))
        if not asset_kind:
            continue
        if idx <= thank_you_index:
            raise ValueError(
                f"Slide {idx + 1} uses a Mermaid diagram before the Thank You slide. Main slides must not use Mermaid; reserve Mermaid diagrams for the appendix after Thank You."
            )
        appendix_diagram_indices.append(idx)

    if appendix_diagram_indices:
        first_appendix_diagram = appendix_diagram_indices[0]
        if not any(
            _is_appendix_section_slide(slide.get('bank_file', ''), slide)
            for slide in slides[thank_you_index + 1:first_appendix_diagram + 1]
        ):
            raise ValueError("Appendix Mermaid diagrams require an Appendix section slide after Thank You and before the first appendix diagram")
        if not (APPENDIX_MIN_DIAGRAMS <= len(appendix_diagram_indices) <= APPENDIX_MAX_DIAGRAMS):
            raise ValueError(
                f"Appendix must contain {APPENDIX_MIN_DIAGRAMS}-{APPENDIX_MAX_DIAGRAMS} Mermaid diagrams; found {len(appendix_diagram_indices)}"
            )
        # Validate that every appendix diagram slide has a text description.
        # Diagram-only slides without explanatory text are not allowed in the appendix.
        for idx in appendix_diagram_indices:
            slide_cfg = slides[idx]
            body = slide_cfg.get('body') or slide_cfg.get('description') or slide_cfg.get('caption') or ''
            if not body.strip():
                raise ValueError(
                    f"Slide {idx + 1} is an appendix diagram slide but has no text description. "
                    f"All appendix diagram slides must include a 'body' (or 'description') field "
                    f"with a text explanation below the diagram."
                )


def _prepare_configured_image_asset(bank_file, image_path, left, top, width, height, diagram_target_class=None):
    resolved = _resolve_asset_path(image_path)
    if not resolved:
        return None

    slot = _image_slots_for_slide(bank_file)[0] if _image_slots_for_slide(bank_file) else None
    slot_class = None
    slot_profile = None
    if slot:
        slot_class, slot_profile = _placeholder_profile_for_slot(slot, bank_file, width=width, height=height)
    asset_kind = _asset_kind_from_path(resolved)
    if asset_kind == 'diagram_source':
        requested_class = _canonical_placeholder_class(diagram_target_class, bank_file=bank_file) if diagram_target_class else None
        if requested_class and slot_class and requested_class != slot_class:
            raise RuntimeError(
                f"{bank_file} targets placeholder class '{slot_class}' but config requested diagram_target_class '{requested_class}'"
            )
        target_class = requested_class or slot_class or 'wide'
        target_profile = PLACEHOLDER_CLASS_PROFILES.get(target_class, PLACEHOLDER_CLASS_PROFILES['wide'])
        source_for_slot = _select_diagram_source_for_slot(resolved, target_class)
        variant_name = f"{source_for_slot.stem}-{target_class}"
        rendered = _render_diagram_source(
            source_for_slot,
            variant_name,
            render_width=target_profile.get('render_width'),
            slot_px=target_profile.get('slot_px'),
        )
        if rendered:
            resolved = rendered
            asset_kind = 'diagram'

    if asset_kind == 'diagram' and resolved.suffix.lower() == '.svg':
        sibling_png = resolved.with_suffix('.png')
        if sibling_png.exists():
            resolved = sibling_png
        else:
            rasterized = convert_svg_to_png(str(resolved), width_px=1800, height_px=1200)
            if rasterized and os.path.exists(rasterized):
                resolved = Path(rasterized)

    if not slot:
        return {
            'path': str(resolved),
            'asset_kind': asset_kind,
            'fit_mode': 'contain' if asset_kind == 'diagram' else 'cover',
            'background': 'transparent',
            'round_corners': asset_kind != 'diagram',
            'corner_radius_px': 8,
        }

    allowed_asset_kinds = slot.get('allowed_asset_kinds') or []
    if asset_kind not in allowed_asset_kinds:
        raise RuntimeError(
            f"{bank_file} placeholder {slot.get('placeholder_idx')} ({slot.get('placeholder_kind')}) does not allow asset kind '{asset_kind}' from {resolved}"
        )

    fit_mode = (slot.get('fit_mode_by_asset_kind') or {}).get(asset_kind, 'contain')
    background = (slot.get('background_by_asset_kind') or {}).get(asset_kind, 'transparent')
    round_corners = bool(slot.get('round_corners', True))
    return {
        'path': str(resolved),
        'asset_kind': asset_kind,
        'fit_mode': fit_mode,
        'background': background,
        'round_corners': round_corners,
        'corner_radius_px': int(slot.get('corner_radius_px', 8)),
    }



AI_TELL_PATTERNS = [
    (re.compile(r'\bAdditionally,\s*', re.IGNORECASE), 'Also, '),
    (re.compile(r'\bFurthermore,\s*', re.IGNORECASE), ''),
    (re.compile(r'\bMoreover,\s*', re.IGNORECASE), ''),
    (re.compile(r'\bIt is important to note that\s*', re.IGNORECASE), ''),
    (re.compile(r'\bIt is worth noting that\s*', re.IGNORECASE), ''),
    (re.compile(r'\bIt should be noted that\s*', re.IGNORECASE), ''),
    (re.compile(r'\bIn conclusion,\s*', re.IGNORECASE), ''),
    (re.compile(r'\bIn summary,\s*', re.IGNORECASE), ''),
    (re.compile(r'\bdelve into\b', re.IGNORECASE), 'explore'),
    (re.compile(r'\bdelve\b', re.IGNORECASE), 'explore'),
    (re.compile(r'\butilize\b', re.IGNORECASE), 'use'),
    (re.compile(r'\butilizes\b', re.IGNORECASE), 'uses'),
    (re.compile(r'\butilizing\b', re.IGNORECASE), 'using'),
    (re.compile(r'\butilization\b', re.IGNORECASE), 'use'),
    (re.compile(r'\bleveraging\b', re.IGNORECASE), 'using'),
    (re.compile(r'\bleverage\b', re.IGNORECASE), 'use'),
    (re.compile(r'\bLeverages\b', re.IGNORECASE), 'Uses'),
    (re.compile(r'\bin order to\b', re.IGNORECASE), 'to'),
    (re.compile(r'\bcommence\b', re.IGNORECASE), 'start'),
    (re.compile(r'\bcommencing\b', re.IGNORECASE), 'starting'),
    (re.compile(r'\bfacilitate\b', re.IGNORECASE), 'enable'),
    (re.compile(r'\bfacilitates\b', re.IGNORECASE), 'enables'),
    (re.compile(r'\bfacilitating\b', re.IGNORECASE), 'enabling'),
    (re.compile(r'\bseamlessly\b', re.IGNORECASE), ''),
    (re.compile(r'\bseamless\b', re.IGNORECASE), 'smooth'),
    (re.compile(r'\bholistic\b', re.IGNORECASE), 'complete'),
    (re.compile(r'\brobust\b', re.IGNORECASE), 'strong'),
    (re.compile(r'\bparadigm\b', re.IGNORECASE), 'model'),
    (re.compile(r'\bsynergy\b', re.IGNORECASE), 'alignment'),
    (re.compile(r'\bsynergies\b', re.IGNORECASE), 'alignments'),
]


def _sanitize_content_text(text):
    """Strip em/en dashes and common AI giveaway patterns from content text."""
    if not isinstance(text, str) or not text.strip():
        return text
    # Replace em dash with comma-space or period depending on context
    result = text.replace('\u2014', ', ')  # em dash → comma space
    result = result.replace('\u2013', '-')  # en dash → hyphen
    result = result.replace('\u2012', '-')  # figure dash → hyphen
    # Clean up double commas or comma-period
    result = result.replace(', ,', ',')
    result = result.replace(',,', ',')
    # Apply AI tell pattern replacements
    for pattern, replacement in AI_TELL_PATTERNS:
        result = pattern.sub(replacement, result)
    # Clean up double spaces
    result = re.sub(r'  +', ' ', result)
    # Clean up leading space after newline
    result = re.sub(r'\n ', '\n', result)
    return result.strip()


def load_config(config_path):
    """Load JSON config."""
    with open(config_path, 'r') as f:
        return json.load(f)


def set_left_align(tf):
    """Legacy helper; retained for compatibility."""
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT


def _is_bullet_line(text):
    stripped = (text or '').lstrip()
    return bool(stripped) and stripped.startswith(BULLET_PREFIXES)


def _strip_bullet_prefix(text):
    stripped = (text or '').lstrip()
    if stripped.startswith(BULLET_PREFIXES):
        return stripped[1:].lstrip()
    return text or ''


def _is_column_title_key(bank_file, key, placeholder_idx=None):
    key_lower = (key or '').lower()
    if key_lower.endswith('_title') and key_lower.startswith('col'):
        return True
    if placeholder_idx is not None:
        return any(spec['placeholder_idx'] == placeholder_idx for spec in _column_title_specs_for_bank(bank_file))
    return False


def _shape_role(bank_file, key, placeholder_idx):
    bank_key = (bank_file or '').lower()
    key_lower = (key or '').lower()
    if placeholder_idx == 0:
        if 'cover' in bank_key:
            return 'cover_title'
        return 'title'
    if key_lower == 'subtitle' and 'cover' in bank_key:
        return 'cover_subtitle'
    if _is_column_title_key(bank_file, key, placeholder_idx):
        return 'mini_title'
    return 'body'


GRID_2X2_PLACEHOLDER_IDXS = {11, 14, 15, 16}
GRID_2X3_PLACEHOLDER_IDXS = {15, 16, 17, 18, 19, 20}
TWO_COL_BODY_PLACEHOLDER_IDXS = {11, 12, 14, 16}
THREE_COL_BODY_PLACEHOLDER_IDXS = {11, 12, 13}
SINGLE_COL_BODY_PLACEHOLDER_IDXS = {11}
SCREENSHOT_BODY_PLACEHOLDER_IDXS = {11}


def _is_grid_text_cell(bank_file, placeholder_idx):
    bank_key = (bank_file or '').lower()
    if ('slide-013' in bank_key or '2x2' in bank_key) and placeholder_idx in GRID_2X2_PLACEHOLDER_IDXS:
        return True
    if ('slide-014' in bank_key or '2x3' in bank_key) and placeholder_idx in GRID_2X3_PLACEHOLDER_IDXS:
        return True
    return False


def _is_label_body_cell(bank_file, placeholder_idx):
    """Detect placeholders that commonly have a label line followed by body text (two-column, three-column, single-column, screenshot layouts)."""
    bank_key = (bank_file or '').lower()
    if ('slide-010' in bank_key or 'two-column' in bank_key or 'two column' in bank_key) and placeholder_idx in TWO_COL_BODY_PLACEHOLDER_IDXS:
        return True
    if ('slide-012' in bank_key or 'three-column' in bank_key or 'three column' in bank_key) and placeholder_idx in THREE_COL_BODY_PLACEHOLDER_IDXS:
        return True
    if ('slide-005' in bank_key or 'single-column' in bank_key or 'single column' in bank_key) and placeholder_idx in SINGLE_COL_BODY_PLACEHOLDER_IDXS:
        return True
    if ('slide-015' in bank_key or 'slide-016' in bank_key or 'screenshot' in bank_key) and placeholder_idx in SCREENSHOT_BODY_PLACEHOLDER_IDXS:
        return True
    return False



LABEL_BODY_WORD_LIMIT = 5  # slightly more generous than COLUMN_TITLE_WORD_LIMIT for labels


def _grid_cell_subtitle_title_and_body(value, bank_file, placeholder_idx):
    is_grid = _is_grid_text_cell(bank_file, placeholder_idx)
    is_label_body = _is_label_body_cell(bank_file, placeholder_idx)
    if not is_grid and not is_label_body:
        return None
    lines = _paragraph_lines(value)
    if len(lines) < 2:
        return None
    title = lines[0].strip()
    if not title or _is_bullet_line(title):
        return None
    word_limit = COLUMN_TITLE_WORD_LIMIT if is_grid else LABEL_BODY_WORD_LIMIT
    if _word_count(title) > word_limit:
        return None
    body_lines = [line for line in lines[1:] if line.strip()]
    if not body_lines:
        return None
    return {'title': title, 'body': body_lines}


def _configure_text_frame(tf, role):
    if role in {'cover_title', 'title', 'mini_title'}:
        # Use NONE for title roles — let fit_text() handle sizing entirely.
        # TEXT_TO_FIT_SHAPE conflicts with fit_text() and can cause horizontal clipping.
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
    else:
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _set_paragraph_layout(paragraph, role, raw_text=''):
    is_bullet = _is_bullet_line(raw_text)
    paragraph.line_spacing = 1.0
    if role in {'title', 'cover_title', 'cover_subtitle', 'mini_title'}:
        paragraph.alignment = PP_ALIGN.LEFT
    else:
        paragraph.alignment = PP_ALIGN.LEFT if is_bullet else PP_ALIGN.JUSTIFY
    p_pr = paragraph._p.get_or_add_pPr()
    for child in list(p_pr):
        if child.tag.endswith('buChar') or child.tag.endswith('buAutoNum') or child.tag.endswith('buBlip') or child.tag.endswith('buNone'):
            p_pr.remove(child)
    p_pr.append(OxmlElement('a:buNone'))
    p_pr.set('marL', '0')
    p_pr.set('indent', '0')


def _suggest_body_max_size(profile_name, estimate=None):
    profile = BODY_FIT_PROFILES.get(profile_name or '', {})
    max_size = profile.get('max_size', 18)
    utilization = (estimate or {}).get('utilization') if estimate else None
    if utilization is not None and utilization < 0.30:
        max_size += 2
    elif utilization is not None and utilization > 0.92:
        max_size -= 1
    return max(12, min(18, int(round(max_size))))


def _fit_text_frame(tf, role, max_size_override=None):
    if role == 'cover_title':
        font_file = str(FIGTREE_BOLD_FONT)
        # Scale max_size based on total text length to prevent single-line overflow.
        # fit_text() only shrinks to fit height, not width — so long single-line titles clip.
        total_chars = sum(len(p.text) for p in tf.paragraphs if p.text)
        if total_chars <= 20:
            max_size = 54
        elif total_chars <= 28:
            max_size = 48
        elif total_chars <= 36:
            max_size = 44
        else:
            max_size = 38
        bold = True
    elif role == 'title':
        font_file = str(FIGTREE_BOLD_FONT)
        # Same treatment for section/slide titles
        total_chars = sum(len(p.text) for p in tf.paragraphs if p.text)
        if total_chars <= 24:
            max_size = 36
        elif total_chars <= 40:
            max_size = 30
        else:
            max_size = 26
        bold = True
    elif role == 'mini_title':
        font_file = str(FIGTREE_BOLD_FONT)
        max_size = 18
        bold = True
    else:
        font_file = str(FIGTREE_FONT)
        max_size = max_size_override or 18
        bold = False
    try:
        tf.fit_text(font_family='Figtree', max_size=max_size, bold=bold, font_file=font_file)
    except Exception:
        pass


# ============================================================================
# TEXT BOX LIMITS - Enforced to prevent text overflow
# ============================================================================

TEXT_LIMITS = {
    ("cover", 0): 60,
    ("section", 0): 60,
    ("separator", 0): 60,
    ("single", 0): 60,
    ("two", 0): 60,
    ("three", 0): 60,
    ("text", 0): 60,
    ("agenda", 0): 60,
    ("table", 0): 60,
    ("screenshot", 0): 60,
    ("picture", 0): 60,
    ("blank", 0): 60,
    ("chart", 0): 60,
    ("cover", 11): 300,
    ("section", 11): 300,
    ("separator", 11): 300,
    ("single", 11): 1400,
    ("two column", 11): 500,
    ("two column", 12): 500,
    ("two column", 14): 500,
    ("two column", 16): 500,
    ("2x2", 11): 250,
    ("2x2", 14): 250,
    ("2x2", 15): 250,
    ("2x2", 16): 250,
    ("screenshot", 11): 500,
    ("agenda", 11): 700,
    ("agenda", 12): 700,
    ("agenda", 13): 700,
    ("agenda", 14): 700,
    ("agenda", 15): 700,
    ("agenda", 16): 700,
    ("agenda", 17): 700,
    ("table", 11): 1000,
}

DEFAULT_LIMIT = 1000


def get_limit(layout_name, placeholder_idx):
    """Get max chars for a given layout and placeholder."""
    layout_lower = layout_name.lower()
    for (pattern, idx), limit in TEXT_LIMITS.items():
        if pattern in layout_lower and idx == placeholder_idx:
            return limit
    return DEFAULT_LIMIT


def truncate_to_limit(text, max_chars):
    """Truncate text to max_chars, preserving paragraph breaks."""
    if not text or not isinstance(text, str):
        return text
    if len(text) <= max_chars:
        return text

    lines = text.split('\n')
    result = []
    total = 0

    for line in lines:
        if total + len(line) + 1 <= max_chars:
            result.append(line)
            total += len(line) + 1
        else:
            remaining = max_chars - total
            if remaining > 10:
                result.append(line[:remaining])
            break

    return '\n'.join(result)


def _paragraph_lines(value):
    normalized = (value or '').replace('\x0b', '\n').replace('\r\n', '\n').replace('\r', '\n')
    return [line.strip() for line in normalized.split('\n') if line.strip()]


def _word_count(text):
    return len([token for token in (text or '').split() if token.strip()])


def _column_title_specs_for_bank(bank_file):
    bank_key = (bank_file or '').replace('.pptx', '').lower()
    for pattern, specs in COLUMN_TITLE_FIELD_SPECS.items():
        if pattern in bank_key:
            return specs
    return []


def _single_line_text(value):
    lines = _paragraph_lines(value)
    return lines[0] if lines else ''


def _emu_to_inches(value):
    return round((value or 0) / EMU_PER_INCH, 3)


def _visible_text_length(text):
    return len((text or '').replace('**', '').strip())


def _infer_fit_profile(bank_file, role, placeholder_idx):
    bank_key = (bank_file or '').lower()
    if role in {'title', 'cover_title', 'cover_subtitle', 'mini_title'}:
        return None
    if 'slide-023' in bank_key or 'full-width' in bank_key:
        return None
    if 'slide-014' in bank_key or '2x3' in bank_key:
        return 'grid_cell_compact'
    if 'slide-013' in bank_key or '2x2' in bank_key:
        return 'grid_cell'
    if 'slide-012' in bank_key or 'three-column' in bank_key or 'three column' in bank_key:
        return 'three_col'
    if 'slide-015' in bank_key or 'slide-016' in bank_key or 'screenshot' in bank_key:
        return 'screenshot_text' if placeholder_idx == 11 else 'body_default'
    if any(token in bank_key for token in ['slide-019', 'slide-020', 'slide-021', 'slide-022', 'one-third', 'wide']):
        return 'picture_narrow_text' if placeholder_idx == 11 else 'body_default'
    if 'slide-010' in bank_key or 'two-column' in bank_key or 'two column' in bank_key:
        return 'body_dense'
    return 'body_default'


def _strip_single_emphasis(text: str) -> str:
    """Remove single-asterisk emphasis (*text*) leaving just the inner text.

    Single * used mid-text for emphasis shows as literal asterisks in PowerPoint
    since we don't map it to italic.  Also cleans up any orphan lone * characters
    that aren't bullet prefixes (those are handled separately by BULLET_PREFIXES).
    """
    # First, strip matched *emphasis* pairs
    result = SINGLE_EMPHASIS_RE.sub(r'\1', text)
    # Then remove any remaining orphan * that aren't part of ** bold markers
    # and aren't at line start (bullet prefix).  Walk through and clean up.
    # Replace lone * that appear mid-text (not preceded/followed by another *)
    cleaned = []
    i = 0
    while i < len(result):
        if result[i] == '*':
            # Check if it's part of a ** bold marker
            if i + 1 < len(result) and result[i + 1] == '*':
                cleaned.append('**')
                i += 2
                continue
            if i > 0 and cleaned and cleaned[-1] == '*':
                # previous was already a lone *, now we see it was ** after all
                # (shouldn't happen after emphasis strip, but be safe)
                i += 1
                continue
            # Lone * at start of line is a bullet prefix — keep it
            line_start = (i == 0) or (i > 0 and result[i - 1] == '\n')
            if line_start:
                cleaned.append('*')
            # else: drop the orphan *
        else:
            cleaned.append(result[i])
        i += 1
    return ''.join(cleaned)


def _normalize_bullet_prefix(text):
    """Convert raw bullet prefixes (*, -, –, —) at line start to a proper • bullet.

    This ensures no literal *, -, –, or — bullet markers leak into the slide.
    Lines already using • are left untouched.
    """
    stripped = (text or '').lstrip()
    if not stripped:
        return text or ''
    for prefix in ('* ', '- ', '– ', '— '):
        if stripped.startswith(prefix):
            return '• ' + stripped[len(prefix):]
    # Handle prefix without trailing space (e.g., "*text")
    for prefix in ('*', '-'):
        if stripped.startswith(prefix) and len(stripped) > 1 and stripped[1] != prefix:
            return '• ' + stripped[1:].lstrip()
    return text or ''


def _parse_body_run_model(text, role, structural_bold=False):
    visible = text or ''
    if role != 'body':
        return [{'text': visible, 'bold': False}]

    # Normalize bullet prefixes to • before any other processing
    visible = _normalize_bullet_prefix(visible)

    if structural_bold and visible.strip():
        return [{'text': visible, 'bold': True}]

    # Strip single-asterisk emphasis before bold parsing to avoid literal * in output
    visible = _strip_single_emphasis(visible)

    if '**' not in visible:
        return [{'text': visible, 'bold': False}]

    parts = BOLD_MARKER_RE.split(visible)
    model = []
    bold_segments = 0
    bold_chars = 0
    total_chars = _visible_text_length(visible)
    for part in parts:
        if not part:
            continue
        is_bold = part.startswith('**') and part.endswith('**') and len(part) > 4
        seg_text = part[2:-2] if is_bold else part
        if not seg_text:
            continue
        if is_bold:
            bold_segments += 1
            bold_chars += len(seg_text.strip())
        model.append({'text': seg_text, 'bold': is_bold})

    if not any(seg['bold'] for seg in model):
        return [{'text': visible.replace('**', ''), 'bold': False}]

    bold_ratio = (bold_chars / total_chars) if total_chars else 0
    if (
        bold_segments > 2
        or bold_ratio > 0.34
        or any(seg['bold'] and len(seg['text'].split()) > 5 for seg in model)
        or all(seg['bold'] for seg in model if seg['text'].strip())
    ):
        # Safety valve: strip all bold markers. Also clean any residual orphan *
        plain = visible.replace('**', '')
        plain = plain.replace('*', '') if plain.count('*') <= 2 else plain
        return [{'text': plain, 'bold': False}]

    return model


def _run_model_text(model):
    return ''.join(segment['text'] for segment in model)


def _estimate_text_utilization(shape, text, profile_name, role='body', bank_file=None, placeholder_idx=None):
    if not profile_name or not text or role != 'body':
        return None
    profile = BODY_FIT_PROFILES[profile_name]
    structural_grid_title = _grid_cell_subtitle_title_and_body(text, bank_file, placeholder_idx)
    structural_title_line = structural_grid_title['title'] if structural_grid_title else None
    width_in = _emu_to_inches(shape.width)
    height_in = _emu_to_inches(shape.height)
    usable_width = max(0.5, width_in * profile['usable_width'])
    usable_height = max(0.4, height_in * profile['usable_height'])
    font_pt = profile['font_pt']
    char_width_in = max(0.03, (font_pt / 72.0) * profile['char_width_em'])
    line_height_in = max(0.10, (font_pt / 72.0) * profile['line_height'])
    chars_per_line = max(8, int(usable_width / char_width_in))
    line_capacity = max(1, int(usable_height / line_height_in))

    consumed_lines = 0.0
    paragraph_models = []
    for line_idx, raw_line in enumerate(_paragraph_lines(text) or ['']):
        model = _parse_body_run_model(
            raw_line,
            role,
            structural_bold=(structural_title_line is not None and line_idx == 0 and raw_line == structural_title_line),
        )
        paragraph_models.append(model)
        visible_line = _run_model_text(model).strip()
        if not visible_line:
            continue
        weighted_chars = 0.0
        for segment in model:
            seg_len = len(segment['text'])
            if segment['bold']:
                seg_len *= 1.08
            weighted_chars += seg_len
        if _is_bullet_line(visible_line):
            weighted_chars += 5
        line_count = max(1, math.ceil(weighted_chars / chars_per_line))
        if _is_bullet_line(visible_line):
            line_count += 0.15
        consumed_lines += line_count

    utilization = round(consumed_lines / line_capacity, 3) if line_capacity else None
    return {
        'profile': profile_name,
        'width_in': width_in,
        'height_in': height_in,
        'chars_per_line': chars_per_line,
        'line_capacity': line_capacity,
        'consumed_lines': round(consumed_lines, 2),
        'utilization': utilization,
        'target_band': profile['target_band'],
    }


def _classify_utilization(estimate):
    if not estimate:
        return {'status': 'not_applicable', 'severity': 'info', 'reason': 'non-body or unprofiled placeholder'}
    profile = BODY_FIT_PROFILES[estimate['profile']]
    utilization = estimate['utilization']
    if utilization is None:
        return {'status': 'not_applicable', 'severity': 'info', 'reason': 'missing utilization score'}
    if utilization < profile['fail_under']:
        return {'status': 'underfilled', 'severity': 'fail', 'reason': f"utilization {utilization:.2f} is far below target band {profile['target_band']}"}
    if utilization < profile['warn_under']:
        return {'status': 'underfilled', 'severity': 'warn', 'reason': f"utilization {utilization:.2f} is below target band {profile['target_band']}"}
    if utilization > profile['fail_over']:
        return {'status': 'overflow-risk', 'severity': 'fail', 'reason': f"utilization {utilization:.2f} exceeds readable limit for {estimate['profile']}"}
    if utilization > profile['warn_over']:
        return {'status': 'overflow-risk', 'severity': 'warn', 'reason': f"utilization {utilization:.2f} is above target band {profile['target_band']}"}
    return {'status': 'ok', 'severity': 'pass', 'reason': f"utilization {utilization:.2f} is inside target band {profile['target_band']}"}


def _save_run_model_after_fit(paragraph, run_model):
    if not any(segment['bold'] for segment in run_model):
        # No bold segments — still sanitize existing runs for em/en dashes
        for run in paragraph.runs:
            run.text = _sanitize_content_text(run.text)
        return
    font_name = None
    font_size = None
    for run in paragraph.runs:
        font_name = font_name or run.font.name
        font_size = font_size or run.font.size
        if font_name and font_size:
            break
    paragraph.clear()
    for segment in run_model:
        run = paragraph.add_run()
        run.text = _sanitize_content_text(segment['text'])
        if font_name:
            run.font.name = font_name
        if font_size:
            run.font.size = font_size
        if segment['bold']:
            run.font.bold = True


def _apply_run_models_after_fit(tf, paragraph_run_models):
    for paragraph, run_model in zip(tf.paragraphs, paragraph_run_models):
        _save_run_model_after_fit(paragraph, run_model)
    # Final pass: sanitize any remaining runs not covered by run models (e.g. title roles)
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            if run.text:
                run.text = _sanitize_content_text(run.text)


def _truncate_title_to_word_limit(title_value, limit=MAIN_TITLE_WORD_LIMIT):
    """Truncate a title to at most `limit` words. Returns the (possibly truncated) single-line title."""
    if not isinstance(title_value, str) or not title_value.strip():
        return title_value
    lines = _paragraph_lines(title_value)
    title_text = lines[0] if lines else title_value.strip()
    words = [w for w in title_text.split() if w]
    if len(words) <= limit:
        return title_text
    return ' '.join(words[:limit])


def _validate_main_title_rules(content):
    title_value = content.get('title')
    if not isinstance(title_value, str) or not title_value.strip():
        return []
    lines = _paragraph_lines(title_value)
    if len(lines) != 1:
        return [f"title must be exactly one paragraph/line (got {len(lines)})"]
    # Word limit is now enforced by auto-truncation upstream; just validate multi-line
    return []


def _validate_column_title_rules(bank_file, content):
    specs = _column_title_specs_for_bank(bank_file)
    violations = []
    for spec in specs:
        title_key = spec['key']
        placeholder_idx = spec['placeholder_idx']
        title_value = content.get(title_key)
        if not isinstance(title_value, str) or not title_value.strip():
            continue

        lines = _paragraph_lines(title_value)
        if len(lines) != 1:
            violations.append(
                f"{title_key} → placeholder idx {placeholder_idx} must be exactly one paragraph/line (got {len(lines)})"
            )
            continue

        line = lines[0]
        if _is_bullet_line(line):
            violations.append(
                f"{title_key} → placeholder idx {placeholder_idx} cannot contain bullets: '{line.strip()}'"
            )
            continue

        wc = _word_count(line)
        if wc > COLUMN_TITLE_WORD_LIMIT:
            violations.append(
                f"{title_key} → placeholder idx {placeholder_idx} has {wc} words (max {COLUMN_TITLE_WORD_LIMIT}): '{line.strip()}'"
            )
    return violations


def _looks_like_asset_path(value):
    if not isinstance(value, str):
        return False
    stripped = value.strip().lower()
    if not stripped:
        return False
    if stripped.startswith(('icon-library/', 'templates/', 'workflow/', 'output/', './', '../', '/')):
        return True
    return bool(re.search(r"\.(svg|png|jpg|jpeg|webp|gif|mmd)$", stripped)) and '/' in stripped


def _validate_agenda_rules(bank_file, content):
    if 'agenda' not in (bank_file or '').lower():
        return []

    violations = []
    agenda_keys = [
        key for key in content.keys()
        if re.fullmatch(r'(?:item|agenda)_\d+', str(key).lower())
    ]
    agenda_indices = sorted({int(str(key).split('_')[1]) for key in agenda_keys})
    if any(index >= AGENDA_ITEM_COUNT for index in agenda_indices):
        violations.append(f"agenda supports at most {AGENDA_ITEM_COUNT} items (item_0..item_6 / agenda_0..agenda_6)")

    for item_index in range(AGENDA_ITEM_COUNT):
        value = _agenda_item_text(content, item_index)
        if not value:
            continue
        lines = _paragraph_lines(value)
        if len(lines) != 1:
            violations.append(f"agenda item {item_index + 1} must be exactly one line / short phrase")
            continue
        line = lines[0].strip()
        if _looks_like_asset_path(line):
            violations.append(f"agenda item {item_index + 1} looks like a file path, not agenda text: '{line}'")
            continue
        if _is_bullet_line(line):
            violations.append(f"agenda item {item_index + 1} cannot contain bullets: '{line}'")
            continue
        wc = _word_count(line)
        if wc > AGENDA_SHORT_PHRASE_WORD_LIMIT:
            violations.append(
                f"agenda item {item_index + 1} has {wc} words (max {AGENDA_SHORT_PHRASE_WORD_LIMIT}): '{line}'"
            )
    return violations


def _validate_hard_text_rules(bank_file, content):
    violations = []
    violations.extend(_validate_main_title_rules(content))
    violations.extend(_validate_column_title_rules(bank_file, content))
    violations.extend(_validate_agenda_rules(bank_file, content))
    return violations


def _canonicalize_column_title_layout_content(bank_file, content):
    normalized = dict(content)
    specs = _column_title_specs_for_bank(bank_file)
    if not specs:
        return normalized

    for spec in specs:
        title_key = spec['key']
        body_key = title_key.replace('_title', '_body')
        fallback_key = spec.get('fallback_key')
        if normalized.get(title_key):
            continue
        fallback_value = normalized.get(fallback_key)
        if not isinstance(fallback_value, str) or not fallback_value.strip():
            continue
        lines = _paragraph_lines(fallback_value)
        if len(lines) >= 2:
            normalized[title_key] = lines[0]
            normalized[body_key] = '\n'.join(lines[1:])
            del normalized[fallback_key]

    return normalized


def _set_paragraph_plain(paragraph):
    _set_paragraph_layout(paragraph, role='body', raw_text=paragraph.text)


def _paragraph_lines_for_role(value, role):
    if role == 'title':
        return [_single_line_text(value) or '']
    if role == 'mini_title':
        return [_single_line_text(value) or '']
    return _paragraph_lines(value) or ['']



def _paragraph_run_models_for_value(value, role, bank_file=None, placeholder_idx=None):
    lines = _paragraph_lines_for_role(value, role)
    structural_grid_title = _grid_cell_subtitle_title_and_body(value, bank_file, placeholder_idx) if role == 'body' else None
    structural_title_line = structural_grid_title['title'] if structural_grid_title else None
    return [
        _parse_body_run_model(
            line,
            role,
            structural_bold=(structural_title_line is not None and i == 0 and line == structural_title_line),
        )
        for i, line in enumerate(lines)
    ]



def _write_text_frame(shape, value, role='body', bank_file=None, placeholder_idx=None):
    tf = shape.text_frame
    tf.clear()
    _configure_text_frame(tf, role)
    if role == 'mini_title':
        lines = _paragraph_lines(value)
        if len(lines) != 1 or _is_bullet_line(lines[0]) or _word_count(lines[0]) > COLUMN_TITLE_WORD_LIMIT:
            raise ValueError(f"Mini-title write rejected for placeholder text: {value!r}")
    lines = _paragraph_lines_for_role(value, role)
    paragraph_run_models = _paragraph_run_models_for_value(value, role, bank_file=bank_file, placeholder_idx=placeholder_idx)

    for i, (line, run_model) in enumerate(zip(lines, paragraph_run_models)):
        paragraph = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        paragraph.text = _run_model_text(run_model)
        _set_paragraph_layout(paragraph, role=role, raw_text=line)

    max_size_override = None
    if role == 'body':
        profile_name = _infer_fit_profile(bank_file, role, placeholder_idx)
        estimate = _estimate_text_utilization(shape, value, profile_name, role=role, bank_file=bank_file, placeholder_idx=placeholder_idx)
        max_size_override = _suggest_body_max_size(profile_name, estimate) if profile_name else None
    _fit_text_frame(tf, role, max_size_override=max_size_override)
    if role == 'body':
        _apply_run_models_after_fit(tf, paragraph_run_models)
    else:
        # Sanitize em/en dashes from title/subtitle runs
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    run.text = _sanitize_content_text(run.text)


def _text_placeholders(placeholders):
    return [
        idx for idx, ph in sorted(placeholders.items())
        if idx not in [0, 10] and hasattr(ph, 'text_frame') and ph.text_frame
        and ph.placeholder_format.type != PP_PLACEHOLDER.TABLE
    ]


def _placeholder_index(bank_file, key, placeholders, used_indices):
    available_indices = list(placeholders.keys())
    text_indices = _text_placeholders(placeholders)
    bank_key = bank_file.replace('.pptx', '').lower()
    key_lower = key.lower()

    if key_lower in ['title', 'heading', 'headline']:
        return 0 if 0 in available_indices else None

    if key_lower == 'subtitle':
        for idx in [11, 12]:
            if idx in available_indices and idx not in used_indices:
                return idx

    if key_lower in ['body', 'content', 'text']:
        return text_indices[0] if text_indices else None

    if 'slide-012' in bank_key:
        explicit = {'col1': 11, 'col2': 12, 'col3': 13}
        if key_lower in explicit and explicit[key_lower] in available_indices:
            return explicit[key_lower]

    if key_lower.startswith('item_') or key_lower.startswith('agenda_'):
        try:
            item_num = int(key_lower.split('_')[1])
        except Exception:
            item_num = None
        if item_num is not None:
            if 'agenda' in bank_key:
                expected = 11 + item_num
                return expected if expected in available_indices else None
            if '2x2' in bank_key:
                ordered = [11, 14, 15, 16]
                return ordered[item_num] if item_num < len(ordered) and ordered[item_num] in available_indices else None
            if '2x3' in bank_key:
                ordered = [15, 16, 17, 18, 19, 20]
                return ordered[item_num] if item_num < len(ordered) and ordered[item_num] in available_indices else None

    if key_lower.startswith('item'):
        suffix = key_lower.replace('item', '')
        digits = ''.join(ch for ch in suffix if ch.isdigit())
        item_num = int(digits) if digits else None
        if item_num is not None:
            if '2x2' in bank_key:
                ordered = [11, 14, 15, 16]
                idx = ordered[item_num - 1] if 1 <= item_num <= len(ordered) else None
                return idx if idx in available_indices else None
            if '2x3' in bank_key:
                ordered = [15, 16, 17, 18, 19, 20]
                idx = ordered[item_num - 1] if 1 <= item_num <= len(ordered) else None
                return idx if idx in available_indices else None

    if key_lower in ['col1', 'col1_title', 'col1_body', 'left', 'left_title', 'left_body']:
        for idx in [11]:
            if idx in available_indices and idx not in used_indices:
                return idx
        return text_indices[0] if text_indices else None

    if key_lower in ['col2', 'col2_title', 'col2_body', 'right', 'right_title', 'right_body']:
        for idx in [12, 14]:
            if idx in available_indices and idx not in used_indices:
                return idx
        if len(text_indices) > 1:
            return text_indices[1]

    if key_lower in ['col3', 'col3_title', 'col3_body']:
        for idx in [13, 15]:
            if idx in available_indices and idx not in used_indices:
                return idx
        if len(text_indices) > 2:
            return text_indices[2]

    for idx in text_indices:
        if idx not in used_indices:
            return idx
    return None


def _add_blank_body_textbox(slide, value):
    # Width 11.8" with 0.77" margins on a 13.33" slide prevents content overflow
    textbox = slide.shapes.add_textbox(Inches(0.77), Inches(1.55), Inches(11.8), Inches(4.9))
    _write_text_frame(textbox, value, bank_file='slide-004-blank.pptx', placeholder_idx=11)
    return textbox


def _populate_table(table, rows):
    num_cols = len(table.columns) if rows else 0
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = _sanitize_content_text(str(val))
            for p in cell.text_frame.paragraphs:
                _set_paragraph_plain(p)
            # Bold header row
            if i == 0:
                for run in cell.text_frame.paragraphs[0].runs:
                    run.font.bold = True
            # Bold first-column labels in data rows (when 3+ columns)
            elif j == 0 and num_cols >= 3:
                for run in cell.text_frame.paragraphs[0].runs:
                    run.font.bold = True


def _fit_table_to_bounds(table, max_width_emu, max_height_emu):
    """Resize table columns and font to fit within placeholder bounds."""
    from pptx.util import Pt

    num_cols = len(table.columns)
    if num_cols > 0 and max_width_emu:
        col_width = max_width_emu // num_cols
        for col in table.columns:
            col.width = col_width

    if not max_height_emu:
        return

    num_rows = len(table.rows)
    chosen_font_pt = 12
    for font_pt in [12, 11, 10, 9, 8, 7]:
        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(font_pt)
        # Approximate row height: font_pt * 1.4 lines * 12700 EMU/pt (compact table rows)
        estimated_row_height = int(font_pt * 1.4 * 12700)
        estimated_total = estimated_row_height * num_rows
        if estimated_total <= max_height_emu:
            chosen_font_pt = font_pt
            break
    else:
        chosen_font_pt = 7

    # Set row heights based on font size only (do NOT stretch to fill placeholder).
    # Stretching causes the last row to overflow when header cells wrap to 2 lines.
    # Use a compact row height with a small padding factor; header row gets 1.8x for wrapped text.
    compact_row_height = int(chosen_font_pt * 1.6 * 12700)  # snug but readable
    header_row_height = int(chosen_font_pt * 2.4 * 12700)    # header may wrap
    for i, row in enumerate(table.rows):
        row.height = header_row_height if i == 0 else compact_row_height

    # Cap rows if total exceeds placeholder height
    min_font_row_height = int(7 * 1.4 * 12700)
    max_rows = max(2, int(max_height_emu / min_font_row_height))
    if num_rows > max_rows:
        # Can't delete rows from python-pptx table easily, but we can clear overflow rows
        for row_idx in range(max_rows, num_rows):
            for cell in table.rows[row_idx].cells:
                cell.text = ""


def build_fit_report(prs, slides_config):
    report = {
        'version': FIT_REPORT_VERSION,
        'slides': [],
        'summary': {'pass': 0, 'warn': 0, 'fail': 0, 'info': 0, 'not_applicable': 0},
    }
    for slide_index, slide in enumerate(prs.slides):
        raw_cfg = slides_config[slide_index] if slide_index < len(slides_config) else {}
        cfg = _canonicalize_column_title_layout_content(raw_cfg.get('bank_file', ''), raw_cfg) if raw_cfg else {}
        placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
        used_indices = set()
        slide_entry = {
            'slide_number': slide_index + 1,
            'bank_file': cfg.get('bank_file'),
            'checks': [],
        }
        for key, value in cfg.items():
            if key in {'bank_file', 'image_path', 'table_data', 'diagram_target_class'} or key.startswith('icon'):
                continue
            if not isinstance(value, str) or not value.strip():
                continue
            idx = _placeholder_index(cfg.get('bank_file', ''), key, placeholders, used_indices)
            if idx is None or idx not in placeholders:
                continue
            shape = placeholders[idx]
            role = _shape_role(cfg.get('bank_file', ''), key, idx)
            profile_name = _infer_fit_profile(cfg.get('bank_file', ''), role, idx)
            estimate = _estimate_text_utilization(
                shape,
                value,
                profile_name,
                role=role,
                bank_file=cfg.get('bank_file', ''),
                placeholder_idx=idx,
            )
            classification = _classify_utilization(estimate)
            check = {
                'key': key,
                'placeholder_idx': idx,
                'role': role,
                'profile': profile_name,
                'status': classification['status'],
                'severity': classification['severity'],
                'reason': classification['reason'],
                'visible_chars': _visible_text_length(value),
                'shape_inches': {'width': _emu_to_inches(shape.width), 'height': _emu_to_inches(shape.height)},
            }
            if estimate:
                check['estimate'] = estimate
            slide_entry['checks'].append(check)
            report['summary'][classification['severity']] = report['summary'].get(classification['severity'], 0) + 1
            used_indices.add(idx)
        report['slides'].append(slide_entry)
    return report


def edit_slide(src_pptx, bank_file, content):
    """Edit slide content using dynamic placeholder discovery and fail-closed coverage."""
    prs = Presentation(src_pptx)
    if not prs.slides:
        return None, []

    slide = prs.slides[0]
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    try:
        layout_name = slide.slide_layout.name if slide.slide_layout else "unknown"
    except Exception:
        layout_name = "unknown"

    expanded_content = _canonicalize_column_title_layout_content(bank_file, content)
    # Auto-truncate main title to word limit (5 words) — ensures single-line rendering
    if isinstance(expanded_content.get('title'), str):
        expanded_content['title'] = _truncate_title_to_word_limit(expanded_content['title'])
    _validate_table_text_requirements(bank_file, expanded_content)
    title_violations = _validate_hard_text_rules(bank_file, expanded_content)
    if title_violations:
        raise ValueError(
            f"Hard text validation failed for {bank_file}: " + '; '.join(title_violations)
        )

    limited_content = {}
    for key, value in expanded_content.items():
        if isinstance(value, str) and value and key not in {'image_path', 'diagram_target_class'}:
            sanitized = _sanitize_content_text(value)
            target_idx = _placeholder_index(bank_file, key, placeholders, set())
            placeholder_idx = 11 if target_idx is None else target_idx
            max_chars = get_limit(layout_name, placeholder_idx)
            limited_content[key] = truncate_to_limit(sanitized, max_chars)
        else:
            limited_content[key] = value

    used_indices = set()
    missing = []

    # Keys reserved for inject_shapes.py post-processing — never touch placeholders
    _RECIPE_SKIP_KEYS = {'recipe'}
    _RECIPE_PREFIXES = (
        'layer', 'step', 'stat', 'left_', 'right_', 'milestone', 'item',
    )
    is_blank_recipe = 'blank' in bank_file.lower() and limited_content.get('recipe')

    for key, value in limited_content.items():
        if key in {'bank_file', 'image_path', 'table_data', 'diagram_target_class'} or str(key).startswith('icon'):
            continue
        # Skip recipe-specific fields on blank recipe slides (handled by inject_shapes.py)
        if is_blank_recipe and (key in _RECIPE_SKIP_KEYS or any(key.startswith(p) for p in _RECIPE_PREFIXES)):
            continue
        if not isinstance(value, str) or not value.strip():
            continue

        idx = _placeholder_index(bank_file, key, placeholders, used_indices)
        if idx is None or idx not in placeholders:
            if key == 'body' and 'blank' in bank_file.lower():
                _add_blank_body_textbox(slide, value)
                continue
            missing.append(key)
            continue

        shape = placeholders[idx]
        if not hasattr(shape, 'text_frame') or not shape.text_frame:
            missing.append(key)
            continue

        role = _shape_role(bank_file, key, idx)
        _write_text_frame(shape, value, role=role, bank_file=bank_file, placeholder_idx=idx)
        used_indices.add(idx)

    if 'table_data' in content and content['table_data']:
        rows = content['table_data']
        table_placeholder = next((ph for ph in slide.placeholders if ph.placeholder_format.type == PP_PLACEHOLDER.TABLE), None)
        if table_placeholder is not None:
            ph_width = table_placeholder.width
            ph_height = table_placeholder.height
            graphic_frame = table_placeholder.insert_table(rows=len(rows), cols=len(rows[0]))
            _populate_table(graphic_frame.table, rows)
            _fit_table_to_bounds(graphic_frame.table, ph_width, ph_height)
        else:
            shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(0.5), Inches(1.5), Inches(9), Inches(3))
            _populate_table(shape.table, rows)
            _fit_table_to_bounds(shape.table, Inches(9), Inches(3))

    return prs, missing


def _sync_slide_content_types(workdir, slide_count):
    """Ensure every merged slide part is registered as a slide content type."""
    content_types_path = os.path.join(workdir, "[Content_Types].xml")
    tree = etree.parse(content_types_path)
    root = tree.getroot()
    ns = '{http://schemas.openxmlformats.org/package/2006/content-types}'
    slide_part_name_prefix = '/ppt/slides/slide'
    slide_content_type = 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml'

    for override in list(root):
        if (
            override.tag == ns + 'Override'
            and (override.get('PartName') or '').startswith(slide_part_name_prefix)
        ):
            root.remove(override)

    existing_part_names = {
        (override.get('PartName') or '')
        for override in root
        if override.tag == ns + 'Override'
    }
    for i in range(1, slide_count + 1):
        part_name = f'{slide_part_name_prefix}{i}.xml'
        if part_name in existing_part_names:
            continue
        override = etree.SubElement(root, ns + 'Override')
        override.set('PartName', part_name)
        override.set('ContentType', slide_content_type)

    tree.write(content_types_path, xml_declaration=True, encoding='UTF-8')



def merge_into_master(edited_dir, slide_count, template_path, output_path):
    """
    Merge edited slides into Master Template.
    Only extract slide XML - avoid duplicate resources.
    """
    workdir = tempfile.mkdtemp()
    
    # Extract Master Template
    with zipfile.ZipFile(template_path, 'r') as z:
        z.extractall(workdir)
    
    # DELETE original slides from Master Template
    slides_dir = os.path.join(workdir, "ppt/slides")
    if os.path.exists(slides_dir):
        for f in os.listdir(slides_dir):
            if f.startswith("slide") and f.endswith(".xml"):
                os.remove(os.path.join(slides_dir, f))
    
    rels_dir = os.path.join(workdir, "ppt/slides/_rels")
    if os.path.exists(rels_dir):
        for f in os.listdir(rels_dir):
            if f.startswith("slide") and f.endswith(".rels"):
                os.remove(os.path.join(rels_dir, f))
    
    # Clean presentation.xml
    pres_xml_path = os.path.join(workdir, "ppt/presentation.xml")
    tree = etree.parse(pres_xml_path)
    root = tree.getroot()
    
    for elem in root.iter():
        if 'sldIdLst' in elem.tag:
            for child in list(elem):
                elem.remove(child)
    
    tree.write(pres_xml_path, xml_declaration=True, encoding='UTF-8')
    
    # Clean relationships
    pres_rels_path = os.path.join(workdir, "ppt/_rels/presentation.xml.rels")
    rels_tree = etree.parse(pres_rels_path)
    rels_root = rels_tree.getroot()
    
    ns = '{http://schemas.openxmlformats.org/package/2006/relationships}'
    for rel in list(rels_root):
        rel_type = rel.get('Type', '')
        if 'relationships/slide' in rel_type and 'slideMaster' not in rel_type:
            rels_root.remove(rel)
    
    # Get max rId
    max_rid = 0
    for rel in rels_root:
        rid = rel.get('Id', '')
        if rid.startswith('rId'):
            try:
                max_rid = max(max_rid, int(rid[3:]))
            except:
                pass
    
    new_rid = max_rid
    
    # Add our slides
    for i in range(slide_count):
        new_rid += 1
        
        # Find edited file for this slide
        edited_files = [f for f in os.listdir(edited_dir) if f.startswith(f"slide-{i:02d}-")]
        if not edited_files:
            continue
        
        edited_path = os.path.join(edited_dir, edited_files[0])
        
        with zipfile.ZipFile(edited_path, 'r') as z:
            # Get slide XML
            slide_files_xml = [f for f in z.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
            if not slide_files_xml:
                continue
            slide_xml = z.read(slide_files_xml[0])
            
            # Get slide relationships
            slide_rels_files = [f for f in z.namelist() if f.startswith('ppt/slides/_rels/') and f.endswith('.rels')]
            slide_rels_xml = b""
            if slide_rels_files:
                slide_rels_xml = z.read(slide_rels_files[0])
            
            # Get media files
            media_files = [f for f in z.namelist() if f.startswith('ppt/media/')]
        
        # Write slide XML
        out_slide_path = os.path.join(workdir, f"ppt/slides/slide{i+1}.xml")
        with open(out_slide_path, 'wb') as f:
            f.write(slide_xml)
        
        # Write slide relationships
        out_rels_path = os.path.join(workdir, f"ppt/slides/_rels/slide{i+1}.xml.rels")
        os.makedirs(os.path.dirname(out_rels_path), exist_ok=True)
        with open(out_rels_path, 'wb') as f:
            f.write(slide_rels_xml)
        
        # Copy media files
        for mf in media_files:
            out_media = os.path.join(workdir, mf)
            if not os.path.exists(out_media):
                os.makedirs(os.path.dirname(out_media), exist_ok=True)
                with open(out_media, 'wb') as f:
                    f.write(z.read(mf))
        
        # Add relationship
        new_rel = etree.SubElement(rels_root, ns + 'Relationship')
        new_rel.set('Id', f'rId{new_rid}')
        new_rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
        new_rel.set('Target', f'slides/slide{i+1}.xml')
        
        # Add to presentation.xml
        for elem in root.iter():
            if 'sldIdLst' in elem.tag:
                new_sldId = etree.SubElement(elem, '{http://schemas.openxmlformats.org/presentationml/2006/main}sldId')
                new_sldId.set('id', str(256 + i))
                new_sldId.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', f'rId{new_rid}')
    
    # Write updated files
    rels_tree.write(pres_rels_path, xml_declaration=True, encoding='UTF-8')
    tree.write(pres_xml_path, xml_declaration=True, encoding='UTF-8')
    _sync_slide_content_types(workdir, slide_count)
    
    # Create output
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as z:
        for root_dir, dirs, files in os.walk(workdir):
            for file in files:
                full_path = os.path.join(root_dir, file)
                arc_path = os.path.relpath(full_path, workdir)
                z.write(full_path, arc_path)
    
    # Cleanup
    shutil.rmtree(workdir)
    
    return output_path




def _icon_file_is_redirect_stub(path):
    try:
        text = Path(path).read_text(encoding='utf-8', errors='ignore')[:256]
    except Exception:
        return False
    return text.startswith('Redirecting to ')


def resolve_icon_path(icon_path):
    """Resolve agenda icons to a real high-quality source asset.

    Priority order:
    1. Template library blue PNG (pre-rendered, brand-correct)
    2. Template library SVG
    3. Lucide SVG (if not a redirect stub)
    4. Original candidate as fallback

    Also strips '.512x512' suffixes from basenames so that e.g.
    'building.512x512.png' can match 'building.png' in the template library.
    """
    candidate = Path(icon_path)
    if not candidate.is_absolute():
        candidate = (ROOT / icon_path).resolve()

    basename = candidate.stem
    # Strip .512x512 (or similar resolution) suffix for template lookup
    import re
    clean_basename = re.sub(r'\.\d+x\d+$', '', basename)

    # Collect candidate basenames (deduplicated, order matters)
    basenames = [basename]
    if clean_basename != basename:
        basenames.append(clean_basename)

    tpl_base = ROOT / 'templates' / 'icon-library'

    # 1. Prefer template blue PNGs (pre-rendered with brand colour)
    for bn in basenames:
        p = (tpl_base / 'png' / 'blue' / f'{bn}.png').resolve()
        if p.exists():
            return str(p)

    # 2. Template SVGs
    for bn in basenames:
        p = (tpl_base / 'svg' / f'{bn}.svg').resolve()
        if p.exists():
            return str(p)

    # 3. Lucide / original SVG if valid
    if candidate.exists() and candidate.suffix.lower() == '.svg' and not _icon_file_is_redirect_stub(candidate):
        return str(candidate)

    # 4. Any existing candidate (e.g. a PNG from lucide)
    if candidate.exists():
        return str(candidate)

    return None


def convert_svg_to_png(svg_path, output_path=None, width_px=512, height_px=512):
    """Convert SVG to a high-resolution PNG via Puppeteer (sole converter).

    Puppeteer is used exclusively so that SVGs containing <foreignObject> or
    complex CSS (e.g. Mermaid diagrams, brand icon SVGs) render correctly.
    cairosvg and rsvg-convert are no longer supported.
    """
    import subprocess
    import shutil
    svg_path = str(svg_path)
    if output_path is None:
        output_path = svg_path.replace('.svg', f'.{width_px}x{height_px}.png')
    # Locate svg_to_png_puppeteer.js relative to this script
    script = Path(__file__).parent / 'svg_to_png_puppeteer.js'
    node_bin = shutil.which('node') or shutil.which('nodejs')
    if not script.exists() or not node_bin:
        return None
    try:
        result = subprocess.run(
            [node_bin, str(script), svg_path, output_path, str(width_px)],
            capture_output=True, text=True, timeout=45,
        )
        if result.returncode == 0 and Path(output_path).exists():
            return output_path
    except Exception:
        pass
    return None


def _prepare_slide_image_asset(image_path, left, top, width, height, *, round_corners=True):
    if not image_path or not os.path.exists(image_path):
        return image_path
    if not round_corners:
        return image_path
    rounded_path = prepare_rounded_image(
        image_path,
        radius_px=8,
        width_emu=width,
        height_emu=height,
        cache_dir=Path(tempfile.gettempdir()) / 'ppt-maker-rounded-images',
    )
    return str(rounded_path) if rounded_path else image_path


def _agenda_item_text(cfg, item_index):
    for key in (f'item_{item_index}', f'agenda_{item_index}'):
        value = cfg.get(key)
        if isinstance(value, str) and value.strip():
            return value.strip()
    return ''


def _remove_shape(shape):
    element = getattr(shape, '_element', None)
    parent = element.getparent() if element is not None else None
    if parent is not None:
        parent.remove(element)


def post_process_pptx(pptx_path, slides_config=None):
    """Apply formatting, inject images/icons, and validate build completeness."""
    from pptx import Presentation

    prs = Presentation(pptx_path)

    if slides_config:
        for i, slide in enumerate(prs.slides):
            if i >= len(slides_config):
                continue
            cfg = slides_config[i]

            if 'agenda' in cfg.get('bank_file', '').lower():
                agenda_placeholders = {
                    shape.placeholder_format.idx: shape
                    for shape in slide.shapes
                    if getattr(shape, 'is_placeholder', False)
                }
                for item_index in range(AGENDA_ITEM_COUNT):
                    text_placeholder_idx = AGENDA_TEXT_PLACEHOLDER_START + item_index
                    icon_placeholder_idx = AGENDA_ICON_PLACEHOLDER_START + item_index
                    text_placeholder = agenda_placeholders.get(text_placeholder_idx)
                    icon_placeholder = agenda_placeholders.get(icon_placeholder_idx)
                    agenda_text = _agenda_item_text(cfg, item_index)

                    if not agenda_text:
                        if text_placeholder is not None:
                            _remove_shape(text_placeholder)
                        if icon_placeholder is not None:
                            _remove_shape(icon_placeholder)
                        continue

                    icon_key = f'icon{item_index + 1}'
                    icon_path = resolve_icon_path(cfg.get(icon_key, '')) if cfg.get(icon_key) else None
                    if icon_placeholder is None or not icon_path or not os.path.exists(icon_path):
                        continue
                    if icon_path.endswith('.svg'):
                        png_path = convert_svg_to_png(icon_path, width_px=512, height_px=512)
                        if png_path:
                            icon_path = png_path
                    if os.path.exists(icon_path):
                        left, top, width, height = icon_placeholder.left, icon_placeholder.top, icon_placeholder.width, icon_placeholder.height
                        _remove_shape(icon_placeholder)
                        add_picture_fit(
                            slide,
                            _prepare_slide_image_asset(icon_path, left, top, width, height, round_corners=False),
                            left,
                            top,
                            width=width,
                            height=height,
                            fit_mode='contain',
                            background='transparent',
                            round_corners=False,
                        )

            img_path = cfg.get('image_path', '')
            if img_path:
                image_slots = _image_slots_for_slide(cfg.get('bank_file', ''))
                preferred_slot = image_slots[0] if image_slots else None
                target_placeholder = None
                if preferred_slot and preferred_slot.get('placeholder_idx') is not None:
                    target_placeholder = next(
                        (
                            shape for shape in slide.shapes
                            if getattr(shape, 'is_placeholder', False)
                            and shape.placeholder_format.idx == preferred_slot.get('placeholder_idx')
                        ),
                        None,
                    )
                if target_placeholder is None:
                    target_placeholder = next(
                        (
                            shape for shape in slide.shapes
                            if getattr(shape, 'is_placeholder', False)
                            and int(shape.placeholder_format.type) in [8, 9, 18]
                        ),
                        None,
                    )
                if target_placeholder is not None:
                    left, top, width, height = target_placeholder.left, target_placeholder.top, target_placeholder.width, target_placeholder.height
                    image_spec = _prepare_configured_image_asset(
                        cfg.get('bank_file', ''),
                        img_path,
                        left,
                        top,
                        width,
                        height,
                        diagram_target_class=cfg.get('diagram_target_class'),
                    )
                    if image_spec and image_spec.get('path'):
                        target_placeholder._element.getparent().remove(target_placeholder._element)
                        add_picture_fit(
                            slide,
                            image_spec['path'],
                            left,
                            top,
                            width=width,
                            height=height,
                            fit_mode=image_spec.get('fit_mode', 'contain'),
                            background=image_spec.get('background', 'transparent'),
                            round_corners=bool(image_spec.get('round_corners', False)),
                            radius_px=int(image_spec.get('corner_radius_px', 8)),
                            cache_dir=Path(tempfile.gettempdir()) / 'ppt-maker-prepared-images',
                        )

    if slides_config:
        for i, slide in enumerate(prs.slides):
            if i >= len(slides_config):
                continue
            cfg = slides_config[i]
            placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
            used_indices = set()
            for key, value in cfg.items():
                if key in {'bank_file', 'image_path', 'table_data', 'diagram_target_class'} or key.startswith('icon'):
                    continue
                if not isinstance(value, str) or not value.strip():
                    continue
                idx = _placeholder_index(cfg.get('bank_file', ''), key, placeholders, used_indices)
                if idx is None or idx not in placeholders:
                    continue
                shape = placeholders[idx]
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    role = _shape_role(cfg.get('bank_file', ''), key, idx)
                    _configure_text_frame(shape.text_frame, role)
                    for paragraph in shape.text_frame.paragraphs:
                        _set_paragraph_layout(paragraph, role=role, raw_text=paragraph.text)
                    profile_name = _infer_fit_profile(cfg.get('bank_file', ''), role, idx)
                    estimate = _estimate_text_utilization(
                        shape,
                        value,
                        profile_name,
                        role=role,
                        bank_file=cfg.get('bank_file', ''),
                        placeholder_idx=idx,
                    ) if profile_name else None
                    _fit_text_frame(shape.text_frame, role, max_size_override=_suggest_body_max_size(profile_name, estimate) if profile_name else None)
                    if role == 'body':
                        _apply_run_models_after_fit(
                            shape.text_frame,
                            _paragraph_run_models_for_value(
                                value,
                                role,
                                bank_file=cfg.get('bank_file', ''),
                                placeholder_idx=idx,
                            ),
                        )
                used_indices.add(idx)
            # NOTE: table cells are already handled by _fit_table_to_bounds in edit_slide.
            # Do NOT re-run fit_text on table cells here — it overrides compact row heights
            # and causes the last row to overflow the placeholder.

    prs.save(pptx_path)

    errors = []
    if slides_config:
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            if i >= len(slides_config):
                continue
            cfg = slides_config[i]
            normalized_cfg = _canonicalize_column_title_layout_content(cfg.get('bank_file', ''), cfg)
            title_violations = _validate_hard_text_rules(cfg.get('bank_file', ''), normalized_cfg)
            if title_violations:
                errors.append(f"slide {i+1}: " + '; '.join(title_violations))
            _RECIPE_SKIP = {'recipe', 'bank_file', 'image_path', 'table_data', 'diagram_target_class'}
            _is_blank_recipe = 'blank' in cfg.get('bank_file', '').lower() and normalized_cfg.get('recipe')
            if _is_blank_recipe:
                # Skip text validation entirely for recipe blank slides - content injected by inject_from_catalog.py
                text_payload = []
            else:
                text_payload = [
                    v for k, v in normalized_cfg.items()
                    if k not in _RECIPE_SKIP
                    and not k.startswith('icon')
                    and isinstance(v, str) and v.strip()
                ]
            visible_text = '\n'.join(
                shape.text for shape in slide.shapes
                if hasattr(shape, 'text') and isinstance(shape.text, str) and shape.text.strip() and 'SettleMint Confidential |' not in shape.text
            )
            for value in text_payload:
                raw_lines = _paragraph_lines(value.replace('**', ''))
                probe_line = _normalize_bullet_prefix(raw_lines[0]) if raw_lines else ''
                probe = probe_line[:24]
                if probe and probe not in visible_text:
                    errors.append(f"slide {i+1}: missing text from config key content starting '{probe}'")
            if cfg.get('table_data') and not any(shape.has_table for shape in slide.shapes):
                errors.append(f"slide {i+1}: expected a native table but none was inserted")
            if cfg.get('image_path') and not any('PICTURE' in str(shape.shape_type) for shape in slide.shapes):
                errors.append(f"slide {i+1}: expected an image/chart fill but no picture shape was found")

    fit_report = None
    # Build set of recipe slide numbers to skip in validation
    recipe_slide_nums = set()
    for si, sc in enumerate(slides_config):
        if sc.get('recipe') and 'blank' in sc.get('bank_file', '').lower():
            recipe_slide_nums.add(si + 1)

    if slides_config:
        prs = Presentation(pptx_path)
        fit_report = build_fit_report(prs, slides_config)
        for slide_entry in fit_report['slides']:
            # Skip validation for recipe blank slides (handled by inject_from_catalog.py)
            if slide_entry['slide_number'] in recipe_slide_nums:
                continue
            for check in slide_entry['checks']:
                if check['severity'] == 'fail':
                    errors.append(
                        f"slide {slide_entry['slide_number']} placeholder {check['placeholder_idx']} ({check['key']}): {check['reason']}"
                    )

    if errors:
        raise RuntimeError('Build validation failed:\n- ' + '\n- '.join(errors))

    return pptx_path, fit_report



def main(config_path):
    """Main build function."""
    # Load config
    config = load_config(config_path)
    
    project = config.get("project", "unnamed")
    output_path = config.get("output", f"output/{project}.pptx")
    slides = config.get("slides", [])
    
    if not slides:
        print("No slides defined in config")
        return

    _validate_mermaid_workflow(slides)
    
    # Setup directories
    script_dir = os.path.dirname(os.path.abspath(__file__))
    bank_dir = os.path.join(os.path.dirname(script_dir), "slide-bank")
    template_path = os.path.join(os.path.dirname(script_dir), "templates/Master Template 2026.pptx")
    
    tmpdir = tempfile.mkdtemp()
    edited_dir = os.path.join(tmpdir, "edited")
    os.makedirs(edited_dir)
    output_path_abs = str((ROOT / output_path).resolve()) if not os.path.isabs(output_path) else output_path
    staged_output_path = os.path.join(tmpdir, Path(output_path_abs).name)
    
    try:
        print(f"Building {project} with {len(slides)} slides...")
        
        # Pre-render all Mermaid diagrams in parallel (if any)
        import time as _time
        _diagram_t0 = _time.monotonic()
        parallel_ok = _prerender_diagrams_parallel(slides)
        _diagram_t1 = _time.monotonic()
        _diagram_elapsed = _diagram_t1 - _diagram_t0
        if _PARALLEL_RENDER_CACHE:
            print(f"  Diagram pre-rendering: {_diagram_elapsed:.1f}s ({len(_PARALLEL_RENDER_CACHE)} diagrams)")
        elif not parallel_ok:
            print(f"  Diagrams will render sequentially during slide assembly")
        
        # Process each slide
        for i, slide_config in enumerate(slides):
            bank_file = slide_config.get("bank_file")
            content = {k: v for k, v in slide_config.items() if k != "bank_file"}
            
            if not bank_file:
                continue
            
            src = os.path.join(bank_dir, bank_file)
            if not os.path.exists(src):
                print(f"  Warning: {src} not found")
                continue
            
            # CRITICAL: Copy to unique name for each slide use
            dst = os.path.join(edited_dir, f"slide-{i:02d}-{bank_file}")
            shutil.copy(src, dst)
            
            # Edit content
            print(f"  {i+1}: {bank_file}")
            prs, missing_keys = edit_slide(dst, bank_file, content)
            if missing_keys:
                # Skip unmapped key errors for blank slides with recipes (handled by inject_from_catalog.py)
                if content.get('recipe') and 'blank' in bank_file.lower():
                    pass  # Recipe fields handled by post-processor
                else:
                    raise RuntimeError(f"Slide {i+1} ({bank_file}) has unmapped content keys: {', '.join(missing_keys)}")
            if prs:
                prs.save(dst)
        
        # Merge into Master Template
        print(f"\nMerging into Master Template...")
        result = merge_into_master(edited_dir, len(slides), template_path, staged_output_path)
        
        # Verify merged slide count before post-processing
        with zipfile.ZipFile(result, 'r') as z:
            slide_files = [f for f in z.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
            print(f"Merged slides: {len(slide_files)}")
        
        # Post-process: apply formatting + fit validation
        print("Applying typography and fit validation...")
        result, fit_report = post_process_pptx(result, slides)

        # Auto-inject recipe shapes for blank slides
        has_recipes = any(
            s.get('recipe') and 'blank' in s.get('bank_file', '').lower()
            for s in slides
        )
        if has_recipes:
            print("Injecting recipe shapes...")
            inject_script = Path(__file__).resolve().parent / 'inject_from_catalog.py'
            if inject_script.exists():
                inject_result = subprocess.run(
                    [sys.executable, str(inject_script), result, config_path],
                    capture_output=True, text=True, timeout=120
                )
                if inject_result.returncode != 0:
                    print(f"  Warning: Recipe injection failed: {inject_result.stderr}")
                else:
                    for line in inject_result.stdout.split('\n'):
                        if line.strip():
                            print(f"  {line.strip()}")
            else:
                print(f"  Warning: inject_from_catalog.py not found at {inject_script}")

            # Validate no empty recipe slides
            from pptx import Presentation as _ValPrs
            val_prs = _ValPrs(result)
            for si, sc in enumerate(slides):
                if sc.get('recipe') and 'blank' in sc.get('bank_file', '').lower():
                    val_slide = val_prs.slides[si]
                    shape_count = len(val_slide.shapes)
                    if shape_count <= 4:
                        raise RuntimeError(
                            f"Slide {si+1} ({sc.get('bank_file')}) uses recipe '{sc['recipe']}' "
                            f"but has only {shape_count} shapes after injection. Recipe injection likely failed."
                        )

        os.makedirs(os.path.dirname(output_path_abs), exist_ok=True)
        shutil.copy2(result, output_path_abs)
        print(f"\nSaved: {output_path_abs}")
        if fit_report:
            report_path = Path(output_path_abs).with_name(f"{Path(output_path_abs).stem}-fit-report.json")
            report_path.write_text(json.dumps(fit_report, indent=2))
            print(f"Fit report: {report_path}")

    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 build_from_config.py <config.json>")
        sys.exit(1)

    main(sys.argv[1])
