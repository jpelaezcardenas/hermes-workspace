
import os
from pptx import Presentation
from pptx.util import Inches

def create_presentation():
    template_path = '/Users/quark/Public/quark/workspace/skills/settlemint-presentation/assets/template.pptx'
    output_path = '/Users/quark/Public/quark/workspace/team-members/navita/output/BIS-Kairos-Tech-Architecture-3slides.pptx'
    diagram_dir = '/Users/quark/Public/quark/workspace/settlemint-office-agents/ppt-maker/output/diagrams'

    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    prs = Presentation(template_path)

    # --- Slide Data ---
    slides_data = [
        {
            "layout_name": "Title - Standard",
            "title": "Technology & Architecture Approach",
            "subtitle": "BIS Project Kairos — SettleMint"
        },
        {
            "layout_name": "Text Left + Image Right",
            "title": "DALP Platform Architecture for Project Kairos",
            "text": "DALP is a composable, EVM-native platform managing the entire tokenised asset lifecycle. For Kairos, it provides token issuance, compliance enforcement (ex-ante, fail-closed), atomic DvP/XvP settlement, identity management, and reconciliation. The SMART Protocol (ERC-3643) ensures every transfer passes modular compliance checks before execution.",
            "image": os.path.join(diagram_dir, "kairos-architecture-stack.png")
        },
        {
            "layout_name": "Text Left + Image Right",
            "title": "Multi-Ledger Environment & Connectivity",
            "text": "Project Kairos spans a 2x2 matrix: EVM/non-EVM crossed with permissioned/permissionless. DALP natively supports both EVM chains (Besu, Ethereum). Solana connectivity uses Hyperlane's cross-chain messaging with Warp Routes and configurable ISMs. Canton integration uses Zenith, SettleMint's EVM compatibility layer, with Daml contracts providing functional equivalence. All deployed within BISIH Azure Cloud.",
            "image": os.path.join(diagram_dir, "kairos-multi-ledger.png")
        },
        {
            "layout_name": "Text Left + Image Right",
            "title": "Cross-Chain Bridge & Settlement Architecture",
            "text": "The bridge uses an escrow-before-burn pattern: tokens are locked on the source ledger before minting on the destination. Two trust models are implemented: a single-operator centralised PoA bridge and a distributed bridge with Hyperlane validator quorum consensus. Cross-chain DvP uses HTLC hashlocks ensuring atomic settlement. Security safeguards include rate limiting, circuit breakers, automated reconciliation, and operation timeouts with recovery queues.",
            "image": os.path.join(diagram_dir, "kairos-bridge-settlement.png")
        },
        {
            "layout_name": "Closing - Thank You",
            "title": "Thank You"
        }
    ]

    # --- Find Layouts ---
    layout_map = {layout.name: layout for layout in prs.slide_layouts}

    # --- Create New Presentation with only the slides we need ---
    new_prs = Presentation()
    new_prs.slide_width = prs.slide_width
    new_prs.slide_height = prs.slide_height

    for data in slides_data:
        layout = layout_map[data["layout_name"]]
        slide = new_prs.slides.add_slide(layout)
        
        # Populate placeholders
        if data.get("title"):
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = data["title"]
        
        if data.get("subtitle"):
            # Find subtitle placeholder by type (often 'CENTER_TITLE' or a body type on title layout)
            for shape in slide.placeholders:
                if shape.placeholder_format.type in ('SUBTITLE', 'BODY') and shape.placeholder_format.idx != 0:
                    shape.text = data["subtitle"]
                    break

        if data.get("text"):
            # Find the main body text placeholder, avoiding the title
            for shape in slide.placeholders:
                if shape.placeholder_format.type in ('BODY', 'OBJECT') and shape.placeholder_format.idx != 0:
                    shape.text = data["text"]
                    break
        
        if data.get("image"):
            # Find the picture placeholder
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 'PICTURE':
                    # This is a more reliable way to add a picture to a picture placeholder
                    placeholder = shape
                    placeholder.insert_picture(data["image"])
                    break

    new_prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
