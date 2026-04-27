#!/usr/bin/env python3
"""Deprecated — use build_from_config.py instead."""
import sys
print("WARNING: build_from_bank.py is deprecated. Use build_from_config.py instead.", file=sys.stderr)
sys.exit(1)

"""
Build PPTX from JSON config using slide bank + Master Template.

KEY PRINCIPLE: Pick slide from bank → fill content in placeholders → merge

This version uses DYNAMIC PLACEHOLDER DISCOVERY to handle all slide types.
"""

import json
import zipfile
import os
import shutil
import tempfile
from lxml import etree
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN


# ============================================================================
# LOAD PLACEHOLDER MAP - Generated from slide bank analysis
# ============================================================================

def load_placeholder_map():
    """Load placeholder structure for each slide template."""
    script_dir = os.path.dirname(os.path.abspath(__file__))
    map_path = os.path.join(script_dir, "placeholder_map.json")
    
    if os.path.exists(map_path):
        with open(map_path, 'r') as f:
            return json.load(f)
    
    # Fallback: analyze on the fly
    return analyze_slide_bank()


def analyze_slide_bank():
    """Analyze slide bank to discover placeholder structures."""
    script_dir = os.path.dirname(os.path.abspath(__file__))
    bank_dir = os.path.join(os.path.dirname(script_dir), "slide-bank")
    
    placeholder_map = {}
    
    for fname in os.listdir(bank_dir):
        if not fname.endswith('.pptx'):
            continue
        
        path = os.path.join(bank_dir, fname)
        try:
            prs = Presentation(path)
            if not prs.slides:
                continue
            
            s = prs.slides[0]
            text_placeholders = []
            
            for ph in s.placeholders:
                idx = ph.placeholder_format.idx
                name = ph.name.lower()
                
                # Skip non-text placeholders
                if 'slide number' in name:
                    continue
                if any(x in name for x in ['online image', 'picture placeholder', 'chart placeholder', 'table placeholder']):
                    continue
                
                text_placeholders.append({'idx': idx, 'name': ph.name})
            
            text_placeholders.sort(key=lambda x: x['idx'])
            placeholder_map[fname.replace('.pptx', '')] = {
                'layout': s.slide_layout.name,
                'placeholders': text_placeholders,
                'count': len(text_placeholders)
            }
        except Exception as e:
            print(f"Warning: Could not analyze {fname}: {e}")
    
    return placeholder_map


# ============================================================================
# SEMANTIC KEY MAPPING
# Maps config keys → placeholder indices based on slide type
# ============================================================================

# Maps semantic keys to expected placeholder positions
# Format: (slide_pattern, key) → list of placeholder indices to try
KEY_MAPPING = {
    # Title - always idx=0
    ('*', 'title'): [0],
    ('*', 'heading'): [0],
    ('*', 'headline'): [0],
    
    # Subtitle / second title
    ('cover', 'subtitle'): [11],
    ('cover', 'subtitle'): [11],
    ('section', 'subtitle'): [11],
    ('section-separator', 'subtitle'): [11],
    
    # Single body content
    ('*', 'body'): [11],
    ('*', 'content'): [11],
    ('*', 'text'): [11],
    
    # Two column
    ('*', 'col1_title'): [11],
    ('*', 'col1'): [11],
    ('*', 'col1_body'): [11],
    ('*', 'left_title'): [11],
    ('*', 'left'): [11],
    ('*', 'left_body'): [11],
    
    ('*', 'col2_title'): [12, 14],
    ('*', 'col2'): [12, 14],
    ('*', 'col2_body'): [12, 14],
    ('*', 'right_title'): [12, 14],
    ('*', 'right'): [12, 14],
    ('*', 'right_body'): [12, 14],
    
    # Three column
    ('*', 'col3_title'): [13, 15],
    ('*', 'col3'): [13, 15],
    ('*', 'col3_body'): [13, 15],
    
    # 2x2 grid items
    ('*', 'item1_title'): [11],
    ('*', 'item1'): [11],
    ('*', 'item1_body'): [11],
    
    ('*', 'item2_title'): [12, 14],
    ('*', 'item2'): [12, 14],
    ('*', 'item2_body'): [12, 14],
    
    ('*', 'item3_title'): [13, 15],
    ('*', 'item3'): [13, 15],
    ('*', 'item3_body'): [13, 15],
    
    ('*', 'item4_title'): [14, 16],
    ('*', 'item4'): [14, 16],
    ('*', 'item4_body'): [14, 16],
    
    # Agenda items (0-6)
    ('*', 'item_0'): [11],
    ('*', 'agenda_0'): [11],
    ('*', 'item_1'): [12],
    ('*', 'agenda_1'): [12],
    ('*', 'item_2'): [13],
    ('*', 'agenda_2'): [13],
    ('*', 'item_3'): [14],
    ('*', 'agenda_3'): [14],
    ('*', 'item_4'): [15],
    ('*', 'agenda_4'): [15],
    ('*', 'item_5'): [16],
    ('*', 'agenda_5'): [16],
    ('*', 'item_6'): [17],
    ('*', 'agenda_6'): [17],
    
    # Generic numbered items (fallback)
    ('*', 'item_7'): [18],
    ('*', 'item_8'): [19],
    ('*', 'item_9'): [20],
}


def get_placeholder_index(bank_file, key, available_indices, used_indices=None):
    """
    Find the correct placeholder index for a given key.
    
    Uses smarter matching: looks at which indices are AVAILABLE (not yet used)
    on this specific slide template.
    
    Args:
        bank_file: e.g., "slide-002-agenda.pptx"
        key: e.g., "title", "item_0", "body"
        available_indices: list of placeholder indices available on this slide
        used_indices: set of indices already used (to avoid duplicates)
    
    Returns:
        The placeholder index to use, or None if not found
    """
    if used_indices is None:
        used_indices = set()
    
    bank_key = bank_file.replace('.pptx', '').lower()
    
    # Get placeholder map for this specific slide
    ph_map = load_placeholder_map()
    slide_key = bank_key
    ph_data = ph_map.get(slide_key, {})
    ph_list = ph_data.get('placeholders', [])
    
    # Extract actual text placeholder indices (excluding 0=title, 10=slide number)
    text_indices = [p['idx'] for p in ph_list if p['idx'] not in [0, 10]]
    
    key_lower = key.lower()
    
    # === TITLE ===
    if key_lower in ['title', 'heading', 'headline']:
        if 0 in available_indices:
            return 0
        return None
    
    # === SUBTITLE (for cover/section slides) ===
    if key_lower == 'subtitle':
        # Try idx=11 first (most common for subtitle)
        for idx in [11, 12]:
            if idx in available_indices and idx not in used_indices:
                return idx
    
    # === BODY / CONTENT (single text block) ===
    if key_lower in ['body', 'content', 'text']:
        for idx in text_indices:
            if idx not in used_indices:
                return idx
    
    # === COLUMN LAYOUTS ===
    # col1 / left column
    if any(k in key_lower for k in ['col1', 'left', '_1']):
        for idx in text_indices:
            if idx not in used_indices:
                return idx
    
    # col2 / right column
    if any(k in key_lower for k in ['col2', 'right', '_2']):
        for idx in text_indices:
            if idx not in used_indices and idx not in [text_indices[0] if text_indices else -1]:
                return idx
    
    # col3 / third column
    if any(k in key_lower for k in ['col3', '_3']):
        for idx in text_indices:
            if idx not in used_indices:
                return idx
    
    # === NUMBERED ITEMS (item_0, item_1, ... or item1, item2, ...) ===
    if key_lower.startswith('item') or key_lower.startswith('agenda'):
        # Extract item number
        item_num = None
        
        # Handle item_0, item_1 format
        if '_' in key_lower:
            try:
                item_num = int(key_lower.split('_')[1])
            except:
                pass
        
        # Handle item1, item2 format
        if item_num is None:
            match = key_lower.replace('item', '').replace('agenda', '')
            try:
                item_num = int(match)
            except:
                pass
        
        if item_num is not None:
            # For agenda slides: item_0 → idx=11, item_1 → idx=12, etc.
            if 'agenda' in bank_key or '002' in bank_key:
                expected = 11 + item_num
                if expected in available_indices:
                    return expected
            
            # For 2x2 grid: use sequential available indices
            # item1 → first available, item2 → second available, etc.
            if '2x2' in bank_key or '013' in bank_key:
                # Skip idx 0 and 10, get body indices
                body_indices = [i for i in text_indices if i != 0 and i != 10]
                if item_num <= len(body_indices):
                    return body_indices[item_num - 1]
            
            # For 2x3 grid
            if '2x3' in bank_key or '014' in bank_key:
                body_indices = [i for i in text_indices if i != 0 and i != 10]
                if item_num <= len(body_indices):
                    return body_indices[item_num - 1]
            
            # For other slides: try 11 + item_num
            expected = 11 + item_num
            if expected in available_indices:
                return expected
    
    # === FALLBACK ===
    # Return first unused text placeholder
    for idx in text_indices:
        if idx not in used_indices:
            return idx
    
    return None


# ============================================================================
# TEXT LIMITS - Prevent overflow
# ============================================================================

TEXT_LIMITS = {
    ('cover', 0): 60,
    ('cover', 11): 300,
    ('section', 0): 60,
    ('section', 11): 300,
    ('agenda', 0): 60,
    ('agenda', 11): 700,
    ('agenda', 12): 700,
    ('agenda', 13): 700,
    ('agenda', 14): 700,
    ('agenda', 15): 700,
    ('agenda', 16): 700,
    ('agenda', 17): 700,
    ('single', 0): 60,
    ('single', 11): 1400,
    ('two', 0): 60,
    ('two', 11): 500,
    ('two', 12): 500,
    ('two', 14): 500,
    ('two', 16): 500,
    ('2x2', 0): 60,
    ('2x2', 11): 250,
    ('2x2', 14): 250,
    ('2x2', 15): 250,
    ('2x2', 16): 250,
    ('2x3', 0): 60,
    ('2x3', 15): 200,
    ('2x3', 16): 200,
    ('2x3', 17): 200,
    ('2x3', 18): 200,
    ('2x3', 19): 200,
    ('2x3', 20): 200,
    ('three', 0): 60,
    ('three', 11): 300,
    ('three', 12): 300,
    ('three', 13): 300,
    ('three', 14): 300,
    ('three', 15): 300,
    ('three', 16): 300,
    ('three', 17): 300,
    ('screenshot', 0): 60,
    ('screenshot', 11): 500,
    ('picture', 0): 60,
    ('picture', 12): 500,
    ('table', 0): 60,
    ('table', 11): 1000,
    ('chart', 0): 60,
    ('chart', 11): 500,
    ('blank', 0): 60,
    ('blank', 11): 2000,
}

DEFAULT_LIMIT = 500


def get_limit(bank_file, placeholder_idx):
    """Get max chars for a given slide type and placeholder."""
    bank_key = bank_file.replace('.pptx', '').lower()
    
    # Extract slide type
    slide_type = bank_key.replace('slide-', '').split('-')[1] if 'slide-' in bank_key else bank_key
    
    # Try exact match
    key = (slide_type, placeholder_idx)
    if key in TEXT_LIMITS:
        return TEXT_LIMITS[key]
    
    # Try just slide type
    for (ptype, pidx), limit in TEXT_LIMITS.items():
        if ptype in slide_type and pidx == placeholder_idx:
            return limit
    
    return DEFAULT_LIMIT


def truncate_to_limit(text, max_chars):
    """Truncate text to max_chars, preserving paragraph breaks."""
    if not text or not isinstance(text, str):
        return text
    if len(text) <= max_chars:
        return text
    
    lines = text.split('\n')
    result = []
    total = 0
    
    for line in lines:
        if total + len(line) + 1 <= max_chars:
            result.append(line)
            total += len(line) + 1
        else:
            remaining = max_chars - total
            if remaining > 10:
                result.append(line[:remaining])
            break
    
    return '\n'.join(result)


# ============================================================================
# SLIDE EDITING
# ============================================================================

def set_left_align(tf):
    """Set left alignment on text frame."""
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT


def edit_slide_dynamic(src_pptx, bank_file, content):
    """
    Edit slide content using DYNAMIC placeholder discovery.
    
    Args:
        src_pptx: path to the bank slide file
        bank_file: filename, e.g., "slide-002-agenda.pptx"
        content: dict of key -> value pairs
    
    Returns:
        Modified Presentation object
    """
    prs = Presentation(src_pptx)
    if not prs.slides:
        return None
    
    slide = prs.slides[0]
    
    # Get available placeholders
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    available_indices = list(placeholders.keys())
    
    # Track which indices we've used
    used_indices = set()
    
    # Apply text limits (don't track used indices yet, just apply limits)
    limited_content = {}
    for key, value in content.items():
        if isinstance(value, str) and value:
            # Find which placeholder this key targets (without tracking used yet)
            target_idx = get_placeholder_index(bank_file, key, available_indices, used_indices=set())
            if target_idx is not None:
                max_chars = get_limit(bank_file, target_idx)
                limited_content[key] = truncate_to_limit(value, max_chars)
            else:
                limited_content[key] = value
        else:
            limited_content[key] = value
    
    # Now apply content to placeholders
    # Priority order: title first, then remaining keys in order
    
    for key, value in limited_content.items():
        if not isinstance(value, str) or not value:
            continue
        
        # Find target placeholder index (track used indices to avoid duplicates)
        target_idx = get_placeholder_index(bank_file, key, available_indices, used_indices)
        
        if target_idx is None:
            # Try fallback: first unused non-title placeholder
            for idx in available_indices:
                if idx != 0 and idx not in used_indices:
                    target_idx = idx
                    break
        
        if target_idx is None or target_idx not in placeholders:
            continue
        
        ph = placeholders[target_idx]
        
        # Skip non-text placeholders
        if not hasattr(ph, 'text_frame') or not ph.text_frame:
            continue
        
        used_indices.add(target_idx)
        
        # Get the text frame
        tf = ph.text_frame
        
        # Clear existing content
        tf.clear()
        
        # Handle multiline text (split by \n)
        lines = value.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line
        
        # Apply left alignment
        set_left_align(tf)
    
    return prs


# ============================================================================
# MERGE INTO MASTER TEMPLATE
# ============================================================================

def _sync_slide_content_types(workdir, slide_count):
    """Ensure every merged slide part is registered as a slide content type."""
    content_types_path = os.path.join(workdir, "[Content_Types].xml")
    tree = etree.parse(content_types_path)
    root = tree.getroot()
    ns = '{http://schemas.openxmlformats.org/package/2006/content-types}'
    slide_part_name_prefix = '/ppt/slides/slide'
    slide_content_type = 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml'

    for override in list(root):
        if (
            override.tag == ns + 'Override'
            and (override.get('PartName') or '').startswith(slide_part_name_prefix)
        ):
            root.remove(override)

    existing_part_names = {
        (override.get('PartName') or '')
        for override in root
        if override.tag == ns + 'Override'
    }
    for i in range(1, slide_count + 1):
        part_name = f'{slide_part_name_prefix}{i}.xml'
        if part_name in existing_part_names:
            continue
        override = etree.SubElement(root, ns + 'Override')
        override.set('PartName', part_name)
        override.set('ContentType', slide_content_type)

    tree.write(content_types_path, xml_declaration=True, encoding='UTF-8')



def merge_into_master(edited_dir, slide_count, template_path, output_path):
    """Merge edited slides into Master Template."""
    workdir = tempfile.mkdtemp()
    
    # Extract Master Template
    with zipfile.ZipFile(template_path, 'r') as z:
        z.extractall(workdir)
    
    # DELETE original slides from Master Template
    slides_dir = os.path.join(workdir, "ppt/slides")
    if os.path.exists(slides_dir):
        for f in os.listdir(slides_dir):
            if f.startswith("slide") and f.endswith(".xml"):
                os.remove(os.path.join(slides_dir, f))
    
    rels_dir = os.path.join(workdir, "ppt/slides/_rels")
    if os.path.exists(rels_dir):
        for f in os.listdir(rels_dir):
            if f.startswith("slide") and f.endswith(".rels"):
                os.remove(os.path.join(rels_dir, f))
    
    # Clean presentation.xml
    pres_xml_path = os.path.join(workdir, "ppt/presentation.xml")
    tree = etree.parse(pres_xml_path)
    root = tree.getroot()
    
    for elem in root.iter():
        if 'sldIdLst' in elem.tag:
            for child in list(elem):
                elem.remove(child)
    
    tree.write(pres_xml_path, xml_declaration=True, encoding='UTF-8')
    
    # Clean relationships
    pres_rels_path = os.path.join(workdir, "ppt/_rels/presentation.xml.rels")
    rels_tree = etree.parse(pres_rels_path)
    rels_root = rels_tree.getroot()
    
    ns = '{http://schemas.openxmlformats.org/package/2006/relationships}'
    for rel in list(rels_root):
        rel_type = rel.get('Type', '')
        if 'relationships/slide' in rel_type and 'slideMaster' not in rel_type:
            rels_root.remove(rel)
    
    # Get max rId
    max_rid = 0
    for rel in rels_root:
        rid = rel.get('Id', '')
        if rid.startswith('rId'):
            try:
                max_rid = max(max_rid, int(rid[3:]))
            except:
                pass
    
    new_rid = max_rid
    
    # Add our slides
    for i in range(slide_count):
        new_rid += 1
        
        # Find edited file for this slide
        edited_files = [f for f in os.listdir(edited_dir) if f.startswith(f"slide-{i:02d}-")]
        if not edited_files:
            continue
        
        edited_path = os.path.join(edited_dir, edited_files[0])
        
        with zipfile.ZipFile(edited_path, 'r') as z:
            # Get slide XML
            slide_files_xml = [f for f in z.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
            if not slide_files_xml:
                continue
            slide_xml = z.read(slide_files_xml[0])
            
            # Get slide relationships
            slide_rels_files = [f for f in z.namelist() if f.startswith('ppt/slides/_rels/') and f.endswith('.rels')]
            slide_rels_xml = b""
            if slide_rels_files:
                slide_rels_xml = z.read(slide_rels_files[0])
            
            # Get media files
            media_files = [f for f in z.namelist() if f.startswith('ppt/media/')]
        
        # Write slide XML
        out_slide_path = os.path.join(workdir, f"ppt/slides/slide{i+1}.xml")
        with open(out_slide_path, 'wb') as f:
            f.write(slide_xml)
        
        # Write slide relationships
        out_rels_path = os.path.join(workdir, f"ppt/slides/_rels/slide{i+1}.xml.rels")
        os.makedirs(os.path.dirname(out_rels_path), exist_ok=True)
        with open(out_rels_path, 'wb') as f:
            f.write(slide_rels_xml)
        
        # Copy media files
        for mf in media_files:
            out_media = os.path.join(workdir, mf)
            if not os.path.exists(out_media):
                os.makedirs(os.path.dirname(out_media), exist_ok=True)
                with open(out_media, 'wb') as f:
                    f.write(z.read(mf))
        
        # Add relationship
        new_rel = etree.SubElement(rels_root, ns + 'Relationship')
        new_rel.set('Id', f'rId{new_rid}')
        new_rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
        new_rel.set('Target', f'slides/slide{i+1}.xml')
        
        # Add to presentation.xml
        for elem in root.iter():
            if 'sldIdLst' in elem.tag:
                new_sldId = etree.SubElement(elem, '{http://schemas.openxmlformats.org/presentationml/2006/main}sldId')
                new_sldId.set('id', str(256 + i))
                new_sldId.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', f'rId{new_rid}')
    
    # Write updated files
    rels_tree.write(pres_rels_path, xml_declaration=True, encoding='UTF-8')
    tree.write(pres_xml_path, xml_declaration=True, encoding='UTF-8')
    _sync_slide_content_types(workdir, slide_count)
    
    # Create output
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as z:
        for root_dir, dirs, files in os.walk(workdir):
            for file in files:
                full_path = os.path.join(root_dir, file)
                arc_path = os.path.relpath(full_path, workdir)
                z.write(full_path, arc_path)
    
    # Cleanup
    shutil.rmtree(workdir)
    
    return output_path


# ============================================================================
# POST PROCESSING
# ============================================================================

def post_process_pptx(pptx_path, slides_config=None):
    """Apply formatting to generated PPTX."""
    from pptx import Presentation
    
    prs = Presentation(pptx_path)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                set_left_align(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            set_left_align(cell.text_frame)
    
    prs.save(pptx_path)
    return pptx_path


# ============================================================================
# MAIN
# ============================================================================

def load_config(config_path):
    """Load JSON config."""
    with open(config_path, 'r') as f:
        return json.load(f)


def main(config_path):
    """Main build function."""
    # Load config
    config = load_config(config_path)
    
    project = config.get("project", "unnamed")
    output_path = config.get("output", f"output/{project}.pptx")
    slides = config.get("slides", [])
    
    if not slides:
        print("No slides defined in config")
        return
    
    # Setup directories
    script_dir = os.path.dirname(os.path.abspath(__file__))
    bank_dir = os.path.join(os.path.dirname(script_dir), "slide-bank")
    template_path = os.path.join(os.path.dirname(script_dir), "templates/Master Template 2026.pptx")
    
    tmpdir = tempfile.mkdtemp()
    edited_dir = os.path.join(tmpdir, "edited")
    os.makedirs(edited_dir)
    
    print(f"Building {project} with {len(slides)} slides...")
    
    # Pre-load placeholder map
    placeholder_map = load_placeholder_map()
    
    # Process each slide
    for i, slide_config in enumerate(slides):
        bank_file = slide_config.get("bank_file")
        content = {k: v for k, v in slide_config.items() if k != "bank_file"}
        
        if not bank_file:
            continue
        
        src = os.path.join(bank_dir, bank_file)
        if not os.path.exists(src):
            print(f"  Warning: {src} not found")
            continue
        
        # Copy to unique name for each slide use
        dst = os.path.join(edited_dir, f"slide-{i:02d}-{bank_file}")
        shutil.copy(src, dst)
        
        # Edit content with dynamic placeholder mapping
        print(f"  {i+1}: {bank_file}")
        prs = edit_slide_dynamic(dst, bank_file, content)
        if prs:
            prs.save(dst)
    
    # Merge into Master Template
    print(f"\nMerging into Master Template...")
    result = merge_into_master(edited_dir, len(slides), template_path, output_path)
    
    # Cleanup
    shutil.rmtree(tmpdir)
    
    print(f"\nSaved: {result}")
    
    # Verify
    with zipfile.ZipFile(result, 'r') as z:
        slide_files = [f for f in z.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
        print(f"Total slides: {len(slide_files)}")
    
    # Post-process
    print("Applying formatting...")
    post_process_pptx(result, slides)
    
    print("✅ Done!")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 build_from_bank.py <config.json>")
        sys.exit(1)
    
    main(sys.argv[1])
